from __future__ import annotations

from pathlib import Path

from harness.preview import build_preview, md_to_html


def test_markdown_preview_escapes_html_and_renders_structures():
    rendered = md_to_html(
        "# Title\n\n<script>alert(1)</script>\n\n"
        "| A | B |\n|---|---|\n| **x** | `y` |\n\n- one\n- two")
    assert "<h1>Title</h1>" in rendered
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in rendered
    assert "<table>" in rendered and "<b>x</b>" in rendered
    assert "<ul><li>one</li><li>two</li></ul>" in rendered


def test_text_json_csv_and_unknown_previews(tmp_path: Path):
    (tmp_path / "data.json").write_text('{"value":"<unsafe>"}', encoding="utf-8")
    (tmp_path / "rows.csv").write_text("name,value\na,1\n", encoding="utf-8")
    (tmp_path / "code.py").write_text("print('<safe>')", encoding="utf-8")
    (tmp_path / "file.bin").write_bytes(b"binary")

    assert "&lt;unsafe&gt;" in build_preview(tmp_path / "data.json")
    csv = build_preview(tmp_path / "rows.csv")
    assert "<th>name</th>" in csv and "<td>1</td>" in csv
    assert "&lt;safe&gt;" in build_preview(tmp_path / "code.py")
    assert "No preview available" in build_preview(tmp_path / "file.bin")


def test_office_and_pdf_previews_render_real_content(tmp_path: Path):
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    import fitz

    doc = Document()
    doc.add_heading("Document heading", level=1)
    doc.add_paragraph("Body text")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header"
    table.cell(1, 0).text = "Cell"
    docx = tmp_path / "sample.docx"
    doc.save(docx)
    doc_html = build_preview(docx)
    assert "Document heading" in doc_html and "<table>" in doc_html

    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Item", "Amount"])
    ws.append(["Rent", 1000])
    ws["B3"] = "=SUM(B2:B2)"
    xlsx = tmp_path / "sample.xlsx"
    wb.save(xlsx)
    xlsx_html = build_preview(xlsx)
    assert "Budget" in xlsx_html and "=SUM(B2:B2)" in xlsx_html

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide title"
    slide.placeholders[1].text = "Slide body"
    pptx = tmp_path / "sample.pptx"
    prs.save(pptx)
    pptx_html = build_preview(pptx)
    assert "Slide title" in pptx_html and "Slide body" in pptx_html

    pdf = tmp_path / "sample.pdf"
    with fitz.open() as document:
        page = document.new_page()
        page.insert_text((72, 72), "PDF text")
        document.save(pdf)
    assert "PDF text" in build_preview(pdf)


def test_broken_preview_returns_readable_error(tmp_path: Path):
    broken = tmp_path / "broken.json"
    broken.write_text("{not json", encoding="utf-8")
    html = build_preview(broken)
    assert "Could not render preview" in html
