"""Reusable native Qt widgets for conversations and job streaming."""
from __future__ import annotations

import base64
import json
import os
import platform
import queue
import re
import sys
import tempfile
import threading
import time
from datetime import datetime
from itertools import islice
from pathlib import Path

from PySide6.QtCore import (
    QEasingCurve,
    QEvent,
    QObject,
    QProcess,
    QPropertyAnimation,
    Qt,
    QThread,
    QTimer,
    QUrl,
    QSize,
    Signal,
)
from PySide6.QtGui import (
    QDesktopServices,
    QDragEnterEvent,
    QDropEvent,
    QIcon,
    QKeyEvent,
    QMouseEvent,
    QPixmap,
    QTextCursor,
)
from PySide6.QtWidgets import (
    QAbstractItemView,
    QDialog,
    QFrame,
    QGraphicsOpacityEffect,
    QGridLayout,
    QHBoxLayout,
    QLabel,
    QLineEdit,
    QListWidget,
    QListWidgetItem,
    QMenu,
    QPlainTextEdit,
    QPushButton,
    QScrollArea,
    QSizePolicy,
    QTextBrowser,
    QTextEdit,
    QToolButton,
    QVBoxLayout,
    QWidget,
)

from .icons import set_svg_icon, svg_icon, svg_pixmap


def clear_layout(layout) -> None:
    while layout.count():
        item = layout.takeAt(0)
        widget = item.widget()
        if widget is not None:
            widget.deleteLater()
        elif item.layout() is not None:
            clear_layout(item.layout())


class IconButton(QPushButton):
    """Accessible button whose chrome is always a bundled SVG."""

    def __init__(self, icon: str, tooltip: str, *, size: int = 16,
                 text: str = "", parent: QWidget | None = None) -> None:
        super().__init__(text, parent)
        self.icon_name = icon
        self.setObjectName("iconButton")
        set_svg_icon(self, icon, size)
        self.setToolTip(tooltip)
        self.setAccessibleName(tooltip)


class SelectButton(QFrame):
    """Predictable cross-platform select control backed by a styled QMenu."""

    changed = Signal(str)

    def __init__(self, *, minimum_width: int = 120,
                 parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self.setObjectName("selectButton")
        self.setCursor(Qt.CursorShape.PointingHandCursor)
        self.setMinimumWidth(minimum_width)
        self._value = ""
        self._options: list[tuple[str, str]] = []
        layout = QHBoxLayout(self)
        layout.setContentsMargins(11, 4, 8, 4)
        layout.setSpacing(7)
        self.label = QLabel()
        self.label.setObjectName("selectLabel")
        self.chevron = QLabel()
        self.chevron.setPixmap(svg_pixmap("chevron_down", 13))
        layout.addWidget(self.label, 1)
        layout.addWidget(self.chevron)

    def set_options(self, options: list[tuple[str, str]],
                    selected: str | None = None) -> None:
        self._options = options
        valid = {value for _label, value in options}
        value = selected if selected in valid else (options[0][1] if options else "")
        self.set_value(value, emit=False)

    def set_value(self, value: str, *, emit: bool = False) -> None:
        self._value = value
        label = next((label for label, item in self._options if item == value), value)
        self.label.setText(label)
        if emit:
            self.changed.emit(value)

    def value(self) -> str:
        return self._value

    def mousePressEvent(self, event: QMouseEvent) -> None:  # noqa: N802
        if event.button() == Qt.MouseButton.LeftButton and self._options:
            menu = QMenu(self)
            menu.setObjectName("selectMenu")
            for label, value in self._options:
                action = menu.addAction(label)
                action.setCheckable(True)
                action.setChecked(value == self._value)
                action.triggered.connect(
                    lambda _checked=False, selected=value: self.set_value(
                        selected, emit=True))
            menu.exec(self.mapToGlobal(self.rect().bottomLeft()))
            event.accept()
            return
        super().mousePressEvent(event)


class TerminalWidget(QWidget):
    """Persistent native terminal for the active workspace."""

    def __init__(self, workspace: Path | None = None) -> None:
        super().__init__()
        self.setObjectName("terminal")
        self.workspace = (workspace or Path.cwd()).resolve()
        layout = QVBoxLayout(self)
        layout.setContentsMargins(8, 8, 8, 8)
        layout.setSpacing(6)
        bar = QHBoxLayout()
        self.location = QLabel(str(self.workspace))
        self.location.setObjectName("terminalLocation")
        clear = IconButton("trash", "Clear terminal", size=14)
        clear.clicked.connect(self.clear_output)
        self.restart_button = IconButton("refresh", "Restart terminal", size=14)
        self.restart_button.clicked.connect(self.restart)
        bar.addWidget(self.location, 1)
        bar.addWidget(clear)
        bar.addWidget(self.restart_button)
        layout.addLayout(bar)
        self.output = QPlainTextEdit()
        self.output.setObjectName("terminalOutput")
        self.output.setReadOnly(True)
        self.output.setMaximumBlockCount(5000)
        layout.addWidget(self.output, 1)
        self.input = QLineEdit()
        self.input.setObjectName("terminalInput")
        self.input.setPlaceholderText("Enter a command…  (↑ ↓ history)")
        self.input.returnPressed.connect(self.submit)
        self.input.installEventFilter(self)
        # Clicking anywhere in the output focuses the prompt, like a terminal.
        self.output.setFocusProxy(self.input)
        self._history: list[str] = []
        self._history_pos: int | None = None
        layout.addWidget(self.input)
        self.process = QProcess(self)
        self.process.setProcessChannelMode(QProcess.ProcessChannelMode.MergedChannels)
        self.process.readyReadStandardOutput.connect(self._read_output)
        self.process.finished.connect(self._finished)
        self.start()

    def eventFilter(self, watched, event) -> bool:  # noqa: N802
        if watched is self.input and event.type() == QEvent.Type.KeyPress:
            if event.key() == Qt.Key.Key_Up and self._history:
                if self._history_pos is None:
                    self._history_pos = len(self._history)
                self._history_pos = max(0, self._history_pos - 1)
                self.input.setText(self._history[self._history_pos])
                return True
            if event.key() == Qt.Key.Key_Down and self._history_pos is not None:
                self._history_pos += 1
                if self._history_pos >= len(self._history):
                    self._history_pos = None
                    self.input.clear()
                else:
                    self.input.setText(self._history[self._history_pos])
                return True
        return super().eventFilter(watched, event)

    def start(self) -> None:
        if self.process.state() != QProcess.ProcessState.NotRunning:
            return
        self.process.setWorkingDirectory(str(self.workspace))
        if platform.system() == "Windows":
            self.process.start("powershell.exe", ["-NoLogo", "-NoProfile"])
        else:
            self.process.start(os.environ.get("SHELL", "/bin/bash"),
                               ["--noprofile", "--norc"])
        if not self.process.waitForStarted(2500):
            self.output.appendPlainText("Could not start the system shell.")

    def restart(self) -> None:
        if self.process.state() != QProcess.ProcessState.NotRunning:
            self.process.kill()
            self.process.waitForFinished(1500)
        self.output.appendPlainText("\n[terminal restarted]")
        self.start()

    def set_workspace(self, workspace: Path) -> None:
        workspace = workspace.resolve()
        if workspace == self.workspace:
            return
        self.workspace = workspace
        self.location.setText(str(workspace))
        if platform.system() == "Windows":
            escaped = str(workspace).replace("'", "''")
            self.write_command(f"Set-Location -LiteralPath '{escaped}'")
        else:
            escaped = str(workspace).replace("'", "'\\''")
            self.write_command(f"cd -- '{escaped}'")

    def submit(self) -> None:
        command = self.input.text().strip()
        if not command:
            return
        if not self._history or self._history[-1] != command:
            self._history.append(command)
        self._history_pos = None
        self.input.clear()
        self.output.appendPlainText(f"> {command}")
        self.write_command(command)

    def write_command(self, command: str) -> None:
        if self.process.state() == QProcess.ProcessState.NotRunning:
            self.start()
        self.process.write((command + "\n").encode("utf-8"))

    def _read_output(self) -> None:
        text = bytes(self.process.readAllStandardOutput().data()).decode(
            "utf-8", errors="replace")
        if text:
            self.output.moveCursor(QTextCursor.MoveOperation.End)
            self.output.insertPlainText(text)
            self.output.moveCursor(QTextCursor.MoveOperation.End)

    def _finished(self, code: int, _status: object = None) -> None:
        self._read_output()
        self.output.appendPlainText(f"\n[shell exited with code {code}]")

    def shutdown(self) -> None:
        if self.process.state() != QProcess.ProcessState.NotRunning:
            self.process.terminate()
            if not self.process.waitForFinished(1000):
                self.process.kill()

    def clear_output(self) -> None:
        self.output.clear()


class _BrowserView(QLabel):
    """Live page surface: clicks, wheel, and keys map onto the real page."""

    clicked_at = Signal(float, float)
    scrolled = Signal(str, int)
    text_typed = Signal(str)
    key_pressed = Signal(str)

    _KEYS: dict[int, str] = {
        Qt.Key.Key_Return: "Enter", Qt.Key.Key_Enter: "Enter",
        Qt.Key.Key_Backspace: "Backspace", Qt.Key.Key_Delete: "Delete",
        Qt.Key.Key_Tab: "Tab", Qt.Key.Key_Escape: "Escape",
        Qt.Key.Key_Up: "ArrowUp", Qt.Key.Key_Down: "ArrowDown",
        Qt.Key.Key_Left: "ArrowLeft", Qt.Key.Key_Right: "ArrowRight",
        Qt.Key.Key_Home: "Home", Qt.Key.Key_End: "End",
        Qt.Key.Key_PageUp: "PageUp", Qt.Key.Key_PageDown: "PageDown",
    }

    def __init__(self) -> None:
        super().__init__(
            "The managed browser has not captured a page yet.\n\n"
            "Enter a URL above — then click, scroll, and type directly "
            "on the page image.")
        self.setObjectName("browserImage")
        self.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.setMinimumHeight(180)
        self.setWordWrap(True)
        self.setFocusPolicy(Qt.FocusPolicy.ClickFocus)
        self.setCursor(Qt.CursorShape.PointingHandCursor)
        self._source = QPixmap()

    def set_source(self, pixmap: QPixmap) -> None:
        self._source = pixmap

    def _map_to_page(self, pos) -> tuple[float, float] | None:
        shown = self.pixmap()
        if self._source.isNull() or shown.isNull():
            return None
        off_x = (self.width() - shown.width()) / 2
        off_y = (self.height() - shown.height()) / 2
        lx, ly = pos.x() - off_x, pos.y() - off_y
        if not (0 <= lx <= shown.width() and 0 <= ly <= shown.height()):
            return None
        return (lx * self._source.width() / shown.width(),
                ly * self._source.height() / shown.height())

    def mousePressEvent(self, event: QMouseEvent) -> None:  # noqa: N802
        if event.button() == Qt.MouseButton.LeftButton:
            mapped = self._map_to_page(event.position())
            if mapped is not None:
                self.setFocus()
                self.clicked_at.emit(*mapped)
                event.accept()
                return
        super().mousePressEvent(event)

    def wheelEvent(self, event) -> None:  # noqa: N802
        delta = event.angleDelta().y()
        if delta:
            self.scrolled.emit("up" if delta > 0 else "down",
                               min(1200, abs(delta) * 2))
            event.accept()
            return
        super().wheelEvent(event)

    def keyPressEvent(self, event: QKeyEvent) -> None:  # noqa: N802
        # Plain int lookup: Qt.Key(...) raises ValueError for keys outside
        # the enum (media keys and other exotics).
        combo = self._KEYS.get(event.key())
        modifiers = event.modifiers()
        control = bool(modifiers & Qt.KeyboardModifier.ControlModifier)
        alt = bool(modifiers & Qt.KeyboardModifier.AltModifier)
        if combo is None and (control or alt) and event.text():
            letter = chr(event.key()) if 32 < event.key() < 127 else ""
            combo = letter.upper() or None
        if combo is not None:
            parts = (["Control"] if control else []) \
                + (["Alt"] if alt else []) \
                + (["Shift"] if modifiers & Qt.KeyboardModifier.ShiftModifier
                   and combo not in {"Enter", "Tab"} else [])
            self.key_pressed.emit("+".join([*parts, combo]))
            event.accept()
            return
        if event.text() and event.text().isprintable():
            self.text_typed.emit(event.text())
            event.accept()
            return
        super().keyPressEvent(event)


def _pty_sequence(event: QKeyEvent) -> str | None:
    """Translate a Qt key event into the byte sequence a terminal expects."""
    key = event.key()
    modifiers = event.modifiers()
    control = bool(modifiers & Qt.KeyboardModifier.ControlModifier)
    shift = bool(modifiers & Qt.KeyboardModifier.ShiftModifier)
    if control and shift:
        return None  # reserved for copy/paste chords handled by the widget
    if control and Qt.Key.Key_A <= key <= Qt.Key.Key_Z:
        return chr(key - Qt.Key.Key_A + 1)
    specials: dict[int, str] = {
        Qt.Key.Key_Return: "\r", Qt.Key.Key_Enter: "\r",
        Qt.Key.Key_Backspace: "\x7f", Qt.Key.Key_Tab: "\t",
        Qt.Key.Key_Backtab: "\x1b[Z", Qt.Key.Key_Escape: "\x1b",
        Qt.Key.Key_Up: "\x1b[A", Qt.Key.Key_Down: "\x1b[B",
        Qt.Key.Key_Right: "\x1b[C", Qt.Key.Key_Left: "\x1b[D",
        Qt.Key.Key_Home: "\x1b[H", Qt.Key.Key_End: "\x1b[F",
        Qt.Key.Key_Delete: "\x1b[3~", Qt.Key.Key_Insert: "\x1b[2~",
        Qt.Key.Key_PageUp: "\x1b[5~", Qt.Key.Key_PageDown: "\x1b[6~",
    }
    if not control and key in specials:
        return specials[key]
    text = event.text()
    if text and text.isprintable():
        return text
    return None


class _WindowsPty:
    def __init__(self, cwd: Path, cols: int, rows: int) -> None:
        from winpty import PtyProcess  # type: ignore[import-not-found]
        self._proc = PtyProcess.spawn(
            "powershell.exe -NoLogo", cwd=str(cwd), dimensions=(rows, cols))

    def read(self) -> str:
        return self._proc.read(4096)

    def write(self, data: str) -> None:
        self._proc.write(data)

    def resize(self, cols: int, rows: int) -> None:
        self._proc.setwinsize(rows, cols)

    def alive(self) -> bool:
        return self._proc.isalive()

    def close(self) -> None:
        try:
            self._proc.terminate(force=True)
        except Exception:
            pass


class _UnixPty:
    """openpty + posix_spawn, never fork(): forking the multithreaded Qt
    process wedges the child on macOS and a wedged child holds pipes (and CI
    steps) open forever."""

    def __init__(self, cwd: Path, cols: int, rows: int) -> None:
        import pty
        import subprocess
        shell = os.environ.get("SHELL", "/bin/bash")
        self._master, slave = pty.openpty()  # type: ignore[attr-defined]
        try:
            self._proc = subprocess.Popen(
                [shell], cwd=str(cwd),
                env={**os.environ, "TERM": "xterm-256color"},
                stdin=slave, stdout=slave, stderr=slave,
                start_new_session=True, close_fds=True)
        except BaseException:
            os.close(self._master)
            raise
        finally:
            os.close(slave)
        self.resize(cols, rows)

    def read(self) -> str:
        data = os.read(self._master, 4096)
        if not data:
            raise EOFError
        return data.decode("utf-8", errors="replace")

    def write(self, data: str) -> None:
        os.write(self._master, data.encode("utf-8"))

    def resize(self, cols: int, rows: int) -> None:
        import fcntl
        import struct
        import termios
        fcntl.ioctl(  # type: ignore[attr-defined]  # unix-only
            self._master, termios.TIOCSWINSZ,  # type: ignore[attr-defined]
            struct.pack("HHHH", rows, cols, 0, 0))

    def alive(self) -> bool:
        return self._proc.poll() is None

    def close(self) -> None:
        import signal
        try:
            os.killpg(self._proc.pid,  # type: ignore[attr-defined]
                      signal.SIGTERM)
        except (OSError, ProcessLookupError):
            pass
        try:
            self._proc.wait(timeout=1.5)
        except Exception:
            try:
                os.killpg(self._proc.pid,  # type: ignore[attr-defined]
                          signal.SIGKILL)  # type: ignore[attr-defined]
            except (OSError, ProcessLookupError):
                pass
        try:
            os.close(self._master)
        except OSError:
            pass


class _TermScreen(QPlainTextEdit):
    """Read-only text surface that forwards raw key sequences to the PTY."""

    keys = Signal(str)

    def __init__(self) -> None:
        super().__init__()
        self.setObjectName("terminalOutput")
        self.setReadOnly(True)
        self.setUndoRedoEnabled(False)
        self.setLineWrapMode(QPlainTextEdit.LineWrapMode.NoWrap)
        self.setFocusPolicy(Qt.FocusPolicy.StrongFocus)
        self.setToolTip("Type directly — Tab completion and shell keys work. "
                        "Ctrl+Shift+C copies, Ctrl+Shift+V pastes.")

    def keyPressEvent(self, event: QKeyEvent) -> None:  # noqa: N802
        modifiers = event.modifiers()
        if (modifiers & Qt.KeyboardModifier.ControlModifier
                and modifiers & Qt.KeyboardModifier.ShiftModifier):
            from PySide6.QtWidgets import QApplication
            if event.key() == Qt.Key.Key_C:
                self.copy()
                event.accept()
                return
            if event.key() == Qt.Key.Key_V:
                text = QApplication.clipboard().text()
                if text:
                    self.keys.emit(text.replace("\r\n", "\r").replace("\n", "\r"))
                event.accept()
                return
        sequence = _pty_sequence(event)
        if sequence is not None:
            self.keys.emit(sequence)
            event.accept()
            return
        super().keyPressEvent(event)


class PtyTerminalWidget(QWidget):
    """A real terminal: ConPTY/pty + pyte emulation, so interactive shells,
    tab completion, history, and TUI redraws behave like a terminal."""

    _data = Signal(str)
    _eof = Signal()

    def __init__(self, workspace: Path | None = None) -> None:
        import pyte
        super().__init__()
        self.setObjectName("terminal")
        self.workspace = (workspace or Path.cwd()).resolve()
        layout = QVBoxLayout(self)
        layout.setContentsMargins(8, 8, 8, 8)
        layout.setSpacing(6)
        bar = QHBoxLayout()
        self.location = QLabel(str(self.workspace))
        self.location.setObjectName("terminalLocation")
        clear = IconButton("trash", "Clear terminal", size=14)
        clear.clicked.connect(self.clear_output)
        self.restart_button = IconButton("refresh", "Restart terminal", size=14)
        self.restart_button.clicked.connect(self.restart)
        bar.addWidget(self.location, 1)
        bar.addWidget(clear)
        bar.addWidget(self.restart_button)
        layout.addLayout(bar)
        self.screen_view = _TermScreen()
        self.input = self.screen_view  # focus-compatibility with the pipe UI
        layout.addWidget(self.screen_view, 1)
        self._cols, self._rows = 100, 28
        self._pyte = pyte.HistoryScreen(
            self._cols, self._rows, history=2000, ratio=0.35)
        self._stream = pyte.Stream(self._pyte)
        self._history_len = 0
        self._history_text: list[str] = []
        self._render_timer = QTimer(self)
        self._render_timer.setSingleShot(True)
        self._render_timer.setInterval(16)
        self._render_timer.timeout.connect(self._render)
        self.screen_view.keys.connect(self._send)
        self._data.connect(self._on_data)
        self._eof.connect(self._on_eof)
        self._pty = self._spawn()
        self._start_reader()
        QTimer.singleShot(0, self._sync_size)

    def _spawn(self):
        backend = _WindowsPty if platform.system() == "Windows" else _UnixPty
        return backend(self.workspace, self._cols, self._rows)

    def _start_reader(self) -> None:
        pty_handle = self._pty

        def read_loop() -> None:
            while True:
                try:
                    chunk = pty_handle.read()
                except (EOFError, OSError, RuntimeError):
                    if pty_handle is self._pty:
                        self._eof.emit()
                    return
                if pty_handle is not self._pty:
                    return
                self._data.emit(chunk)

        threading.Thread(target=read_loop, name="lmh-terminal-pty",
                         daemon=True).start()

    def _on_data(self, chunk: str) -> None:
        self._stream.feed(chunk)
        if not self._render_timer.isActive():
            self._render_timer.start()

    def _on_eof(self) -> None:
        self._stream.feed("\r\n[shell exited — restart to reopen]\r\n")
        self._render()

    def _render(self) -> None:
        history = self._pyte.history.top
        if len(history) != self._history_len:
            self._history_text = [
                "".join(line[x].data for x in range(self._cols)).rstrip()
                for line in history]
            self._history_len = len(history)
        display = [line.rstrip() for line in self._pyte.display]
        keep = self._pyte.cursor.y + 1
        while len(display) > keep and not display[-1]:
            display.pop()
        bar = self.screen_view.verticalScrollBar()
        stick = bar.maximum() - bar.value() < 40
        self.screen_view.setPlainText(
            "\n".join([*self._history_text, *display]))
        if stick:
            bar.setValue(bar.maximum())
        cursor = self.screen_view.textCursor()
        block = self.screen_view.document().findBlockByNumber(
            len(self._history_text) + self._pyte.cursor.y)
        if block.isValid():
            position = block.position() + min(self._pyte.cursor.x,
                                              max(0, block.length() - 1))
            cursor.setPosition(position)
            selection = QTextEdit.ExtraSelection()
            selection.cursor = cursor
            selection.cursor.movePosition(
                QTextCursor.MoveOperation.NextCharacter,
                QTextCursor.MoveMode.KeepAnchor)
            selection.format.setBackground(self.palette().highlight())
            self.screen_view.setExtraSelections([selection])

    def _send(self, data: str) -> None:
        if self._pty.alive():
            try:
                self._pty.write(data)
            except OSError:
                pass

    def _sync_size(self) -> None:
        from PySide6.QtGui import QFontMetricsF
        metrics = QFontMetricsF(self.screen_view.font())
        char_width = max(4.0, metrics.horizontalAdvance("M"))
        line_height = max(8.0, metrics.lineSpacing())
        viewport = self.screen_view.viewport().size()
        cols = max(20, int(viewport.width() / char_width) - 1)
        rows = max(5, int(viewport.height() / line_height))
        if (cols, rows) == (self._cols, self._rows):
            return
        self._cols, self._rows = cols, rows
        try:
            self._pyte.resize(rows, cols)
            self._pty.resize(cols, rows)
        except Exception:
            return
        self._history_len = -1  # re-render history at the new width
        self._render()

    def resizeEvent(self, event) -> None:  # noqa: N802
        super().resizeEvent(event)
        QTimer.singleShot(60, self._sync_size)

    def set_workspace(self, workspace: Path) -> None:
        workspace = workspace.resolve()
        if workspace == self.workspace:
            return
        self.workspace = workspace
        self.location.setText(str(workspace))
        if platform.system() == "Windows":
            escaped = str(workspace).replace("'", "''")
            self._send(f"Set-Location -LiteralPath '{escaped}'\r")
        else:
            escaped = str(workspace).replace("'", "'\\''")
            self._send(f"cd -- '{escaped}'\r")

    def restart(self) -> None:
        old = self._pty
        try:
            replacement = self._spawn()
        except Exception as exc:
            self._stream.feed(f"\r\n[could not restart the shell: {exc}]\r\n")
            self._render()
            return
        # Swap before closing: the old reader thread compares identity and
        # exits silently instead of printing "[shell exited]" post-restart.
        self._pty = replacement
        old.close()
        self.clear_output()
        self._start_reader()
        self._sync_size()

    def clear_output(self) -> None:
        self._pyte.history.top.clear()
        self._pyte.history.bottom.clear()
        self._history_len = 0
        self._history_text = []
        self._pyte.reset()
        self._render()

    def shutdown(self) -> None:
        self._pty.close()


def create_terminal(
        workspace: Path | None = None) -> "PtyTerminalWidget | TerminalWidget":
    """A real PTY terminal when the runtime supports it, else the pipe UI.

    LMH_PIPE_TERMINAL forces the pipe fallback — packaged smoke tests use it
    so release verification never depends on spawning an interactive shell.
    """
    if os.environ.get("LMH_PIPE_TERMINAL"):
        return TerminalWidget(workspace)
    try:
        return PtyTerminalWidget(workspace)
    except Exception:
        return TerminalWidget(workspace)


class BrowserPanel(QWidget):
    """Native inspector/control surface for the managed Playwright browser."""

    result_ready = Signal(str, bytes)

    def __init__(self, workspace: Path | None = None) -> None:
        super().__init__()
        self.setObjectName("browserPanel")
        self.workspace = (workspace or Path.cwd()).resolve()
        self._busy = False
        self._pixmap = QPixmap()
        self._type_buffer = ""
        self._type_timer = QTimer(self)
        self._type_timer.setSingleShot(True)
        self._type_timer.setInterval(350)
        self._type_timer.timeout.connect(self._flush_typed)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(8, 8, 8, 8)
        layout.setSpacing(6)
        bar = QHBoxLayout()
        back = IconButton("arrow_left", "Back", size=14)
        forward = IconButton("arrow_right", "Forward", size=14)
        reload_button = IconButton("refresh", "Reload", size=14)
        back.clicked.connect(lambda: self.run_action("back"))
        forward.clicked.connect(lambda: self.run_action("forward"))
        reload_button.clicked.connect(lambda: self.run_action("reload"))
        self.address = QLineEdit()
        self.address.setPlaceholderText("Enter a public URL…")
        self.address.returnPressed.connect(self.navigate)
        go = IconButton("arrow_right", "Navigate", size=14)
        go.clicked.connect(self.navigate)
        bar.addWidget(back)
        bar.addWidget(forward)
        bar.addWidget(reload_button)
        bar.addWidget(self.address, 1)
        bar.addWidget(go)
        layout.addLayout(bar)
        self.image = _BrowserView()
        self.image.clicked_at.connect(
            lambda x, y: self.run_action("click_at", x=x, y=y))
        self.image.scrolled.connect(
            lambda direction, amount: self.run_action(
                "scroll", direction=direction, amount=amount))
        self.image.text_typed.connect(self._buffer_typed)
        self.image.key_pressed.connect(
            lambda key: self.run_action("press", key=key))
        layout.addWidget(self.image, 2)
        self.state = QPlainTextEdit()
        self.state.setObjectName("browserState")
        self.state.setReadOnly(True)
        self.state.setMaximumHeight(220)
        layout.addWidget(self.state, 1)
        self.result_ready.connect(self._show_result)

    def _buffer_typed(self, text: str) -> None:
        self._type_buffer += text
        self._type_timer.start()

    def _flush_typed(self) -> None:
        if self._busy:
            self._type_timer.start()  # retry once the current action lands
            return
        buffer, self._type_buffer = self._type_buffer, ""
        if buffer:
            self.run_action("type_text", text=buffer)

    def set_workspace(self, workspace: Path) -> None:
        self.workspace = workspace.resolve()

    def navigate(self) -> None:
        url = self.address.text().strip()
        if url:
            self.run_action("open", url=url)

    def run_action(self, action: str, **kwargs) -> None:
        if self._busy:
            return
        self._busy = True
        self.state.setPlainText(f"{action.title()}…")

        def work() -> None:
            try:
                from ..browser import control
                result = control(action, screenshot_dir=self.workspace, **kwargs)
                report, image = self._decode_result(result)
            except Exception as exc:
                report = f"Error: {type(exc).__name__}: {exc}"
                image = b""
            self.result_ready.emit(report, image)

        threading.Thread(target=work, name="lmh-browser-panel", daemon=True).start()

    @staticmethod
    def _decode_result(result: str) -> tuple[str, bytes]:
        if not result.startswith("__MCP_IMAGE_RESULT__:"):
            return result, b""
        try:
            payload = json.loads(result.split(":", 1)[1])
            report = str(payload.get("report", ""))
            images = payload.get("images") or []
            raw = base64.b64decode(images[0]["data"], validate=True) if images else b""
            return report, raw
        except (json.JSONDecodeError, KeyError, TypeError, ValueError):
            return "Browser returned malformed state.", b""

    def _show_result(self, report: str, image: bytes) -> None:
        self._busy = False
        self.state.setPlainText(report)
        url = next((line[5:].strip() for line in report.splitlines()
                    if line.startswith("URL: ")), "")
        if url:
            self.address.setText(url)
        pixmap = QPixmap()
        if image and pixmap.loadFromData(image):
            self._pixmap = pixmap
            self.image.set_source(pixmap)
            self._fit_image()
        elif report.startswith("Error"):
            self.image.setText(report)

    def _fit_image(self) -> None:
        if not self._pixmap.isNull():
            self.image.setPixmap(self._pixmap.scaled(
                max(120, self.image.width() - 8),
                max(120, self.image.height() - 8),
                Qt.AspectRatioMode.KeepAspectRatio,
                Qt.TransformationMode.SmoothTransformation))

    def resizeEvent(self, event) -> None:  # noqa: N802
        super().resizeEvent(event)
        self._fit_image()


class NativeEmbeddedBrowserPanel(QWidget):
    """The agent's real Chromium, embedded live in the panel (Windows).

    The shared Playwright browser runs headful; this panel captures its
    top-level window with createWindowContainer, so the user browses the
    actual page with real input while the agent drives the same browser.
    While a job runs, input to the embedded window is disabled.
    """

    def __init__(self, workspace: Path | None = None) -> None:
        super().__init__()
        self.setObjectName("browserPanel")
        self.workspace = (workspace or Path.cwd()).resolve()
        self._hwnd: int | None = None
        self._busy_lock = False
        layout = QVBoxLayout(self)
        layout.setContentsMargins(8, 8, 8, 8)
        layout.setSpacing(6)
        bar = QHBoxLayout()
        back = IconButton("arrow_left", "Back", size=14)
        forward = IconButton("arrow_right", "Forward", size=14)
        reload_button = IconButton("refresh", "Reload", size=14)
        back.clicked.connect(lambda: self._drive("back"))
        forward.clicked.connect(lambda: self._drive("forward"))
        reload_button.clicked.connect(lambda: self._drive("reload"))
        self.address = QLineEdit()
        self.address.setPlaceholderText("Enter a URL…")
        self.address.returnPressed.connect(self.navigate)
        go = IconButton("arrow_right", "Navigate", size=14)
        go.clicked.connect(self.navigate)
        self.lock_hint = QLabel("Assistant is controlling the browser")
        self.lock_hint.setObjectName("browserLockText")
        self.lock_hint.hide()
        bar.addWidget(back)
        bar.addWidget(forward)
        bar.addWidget(reload_button)
        bar.addWidget(self.address, 1)
        bar.addWidget(go)
        layout.addLayout(bar)
        layout.addWidget(self.lock_hint)
        self.placeholder = QLabel(
            "The live browser appears here the first time it is used.\n\n"
            "Enter a URL above to start it — you can then browse directly, "
            "and watch the assistant work in the same window.")
        self.placeholder.setObjectName("browserImage")
        self.placeholder.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.placeholder.setWordWrap(True)
        self.placeholder.setMinimumHeight(180)
        layout.addWidget(self.placeholder, 1)
        self._layout = layout
        # Capture the browser window whenever it exists — whether the user
        # started it from the address bar or the agent launched it mid-task.
        self._capture_timer = QTimer(self)
        self._capture_timer.setInterval(1500)
        self._capture_timer.timeout.connect(self._try_capture)
        self._capture_timer.start()

    def _drive(self, action: str, **kwargs) -> None:
        def work() -> None:
            from .. import browser
            browser.control(action, screenshot_dir=self.workspace, **kwargs)

        threading.Thread(target=work, name="lmh-embed-browse",
                         daemon=True).start()

    def navigate(self) -> None:
        raw = self.address.text().strip()
        if not raw:
            return
        if not re.match(r"^https?://", raw, re.IGNORECASE):
            raw = "https://" + raw
        self._drive("open", url=raw)

    def _try_capture(self) -> None:
        if self._hwnd is not None:
            return
        from .. import browser
        hwnd = browser.embedded_window_handle()
        if not hwnd:
            return
        from PySide6.QtGui import QWindow
        foreign = QWindow.fromWinId(hwnd)
        if foreign is None:
            return
        container = QWidget.createWindowContainer(foreign, self)
        container.setMinimumHeight(180)
        self._layout.replaceWidget(self.placeholder, container)
        self.placeholder.hide()
        self.placeholder.deleteLater()
        self._layout.setStretchFactor(container, 1)
        self._hwnd = hwnd
        self._apply_lock()

    def set_locked(self, locked: bool) -> None:
        self._busy_lock = locked
        self._apply_lock()

    def _apply_lock(self) -> None:
        self.address.setDisabled(self._busy_lock)
        self.lock_hint.setVisible(self._busy_lock and self._hwnd is not None)
        if self._hwnd is not None and sys.platform == "win32":
            import ctypes
            ctypes.windll.user32.EnableWindow(
                self._hwnd, 0 if self._busy_lock else 1)

    def set_workspace(self, workspace: Path) -> None:
        self.workspace = workspace.resolve()


def create_browser_panel(
        workspace: Path | None = None,
) -> "NativeEmbeddedBrowserPanel | BrowserPanel":
    """Live embedded Chromium on Windows; the screenshot panel elsewhere."""
    if os.environ.get("LMH_HEADFUL_BROWSER") == "1" \
            and platform.system() == "Windows":
        try:
            return NativeEmbeddedBrowserPanel(workspace)
        except Exception:
            pass
    return BrowserPanel(workspace)


class ArtifactPreview(QWidget):
    """Native preview for workspace artifacts without an embedded web app."""

    def __init__(self) -> None:
        super().__init__()
        self.path: Path | None = None
        root = QVBoxLayout(self)
        root.setContentsMargins(8, 8, 8, 8)
        root.setSpacing(6)
        bar = QHBoxLayout()
        self.title = QLabel("Select a workspace file to preview")
        self.title.setObjectName("artifactTitle")
        self.open_button = IconButton("download", "Open with system app", size=14)
        self.open_button.setEnabled(False)
        self.open_button.clicked.connect(self.open_external)
        bar.addWidget(self.title, 1)
        bar.addWidget(self.open_button)
        root.addLayout(bar)
        self.body = QWidget()
        self.body_layout = QVBoxLayout(self.body)
        self.body_layout.setContentsMargins(0, 0, 0, 0)
        self.body_layout.addStretch(1)
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.scroll_area.setWidget(self.body)
        root.addWidget(self.scroll_area, 1)

    def preview(self, path: Path) -> None:
        self.path = path.resolve()
        self.title.setText(path.name)
        self.open_button.setEnabled(True)
        clear_layout(self.body_layout)
        try:
            suffix = path.suffix.lower()
            if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}:
                self._preview_image(path)
            elif suffix == ".pdf":
                self._preview_pdf(path)
            elif suffix == ".docx":
                self._preview_text(self._docx_text(path))
            elif suffix in {".xlsx", ".xlsm"}:
                self._preview_text(self._xlsx_text(path))
            elif suffix == ".pptx":
                self._preview_text(self._pptx_text(path))
            elif path.stat().st_size <= 2 * 1024 * 1024:
                self._preview_text(path.read_text(encoding="utf-8", errors="replace"))
            else:
                self._preview_text(
                    "This file is too large for an inline text preview. Open it "
                    "with the system application instead.")
        except Exception as exc:
            self._preview_text(f"Could not preview {path.name}: {type(exc).__name__}: {exc}")
        self.body_layout.addStretch(1)

    def _preview_image(self, path: Path) -> None:
        pixmap = QPixmap(str(path))
        if pixmap.isNull():
            raise ValueError("image decoder rejected the file")
        image = QLabel()
        image.setAlignment(Qt.AlignmentFlag.AlignTop | Qt.AlignmentFlag.AlignHCenter)
        image.setPixmap(pixmap.scaled(
            1100, 900, Qt.AspectRatioMode.KeepAspectRatio,
            Qt.TransformationMode.SmoothTransformation))
        self.body_layout.addWidget(image)

    def _preview_pdf(self, path: Path) -> None:
        import fitz  # type: ignore[import-untyped]
        with fitz.open(str(path)) as document:
            for page_number in range(min(document.page_count, 8)):
                pix = document[page_number].get_pixmap(dpi=105)
                preview = QPixmap()
                preview.loadFromData(pix.tobytes("png"))
                label = QLabel()
                label.setAlignment(Qt.AlignmentFlag.AlignHCenter)
                label.setPixmap(preview)
                label.setToolTip(f"Page {page_number + 1}")
                self.body_layout.addWidget(label)
            if document.page_count > 8:
                self.body_layout.addWidget(QLabel(
                    f"Showing 8 of {document.page_count} pages. Open the file for the rest."))

    def _preview_text(self, text: str) -> None:
        view = QPlainTextEdit()
        view.setObjectName("artifactText")
        view.setReadOnly(True)
        view.setPlainText(text[:200_000])
        view.setMinimumHeight(500)
        self.body_layout.addWidget(view)

    @staticmethod
    def _docx_text(path: Path) -> str:
        from docx import Document
        document = Document(str(path))
        lines = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            lines.extend(" | ".join(cell.text for cell in row.cells)
                         for row in table.rows)
        return "\n".join(lines)

    @staticmethod
    def _xlsx_text(path: Path) -> str:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
        workbook = load_workbook(path, read_only=True, data_only=False)
        lines = []
        for sheet in workbook.worksheets:
            lines.append(f"[{sheet.title}]")
            for row in islice(sheet.iter_rows(values_only=True), 150):
                lines.append("\t".join("" if value is None else str(value)
                                       for value in row[:40]))
        workbook.close()
        return "\n".join(lines)

    @staticmethod
    def _pptx_text(path: Path) -> str:
        from pptx import Presentation
        presentation = Presentation(str(path))
        lines = []
        for index, slide in enumerate(presentation.slides, 1):
            lines.append(f"[Slide {index}]")
            lines.extend(str(shape.text) for shape in slide.shapes
                         if hasattr(shape, "text") and shape.text)
        return "\n".join(lines)

    def open_external(self) -> None:
        if self.path is not None:
            QDesktopServices.openUrl(QUrl.fromLocalFile(str(self.path)))


class ComposerEdit(QPlainTextEdit):
    submitted = Signal()
    stop_requested = Signal()
    files_dropped = Signal(list)

    def __init__(self, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self.setAcceptDrops(True)

    def keyPressEvent(self, event: QKeyEvent) -> None:  # noqa: N802
        if event.key() == Qt.Key.Key_Escape:
            self.stop_requested.emit()
            event.accept()
            return
        if (event.key() in {Qt.Key.Key_Return, Qt.Key.Key_Enter}
                and not event.modifiers() & Qt.KeyboardModifier.ShiftModifier):
            self.submitted.emit()
            event.accept()
            return
        if (event.key() == Qt.Key.Key_V
                and event.modifiers() & Qt.KeyboardModifier.ControlModifier):
            from PySide6.QtWidgets import QApplication
            mime = QApplication.clipboard().mimeData()
            if mime.hasImage():
                image = QApplication.clipboard().image()
                path = Path(tempfile.gettempdir()) / (
                    f"little-harness-paste-{time.time_ns()}.png")
                if image.save(str(path), b"PNG"):
                    self.files_dropped.emit([str(path)])
                    event.accept()
                    return
            if mime.hasUrls():
                paths = [url.toLocalFile() for url in mime.urls()
                         if url.isLocalFile()]
                if paths:
                    self.files_dropped.emit(paths)
                    event.accept()
                    return
        super().keyPressEvent(event)

    def dragEnterEvent(self, event: QDragEnterEvent) -> None:  # noqa: N802
        if event.mimeData().hasUrls() and any(
                url.isLocalFile() for url in event.mimeData().urls()):
            event.acceptProposedAction()
            return
        super().dragEnterEvent(event)

    def dropEvent(self, event: QDropEvent) -> None:  # noqa: N802
        paths = [url.toLocalFile() for url in event.mimeData().urls()
                 if url.isLocalFile()]
        if paths:
            self.files_dropped.emit(paths)
            event.acceptProposedAction()
            return
        super().dropEvent(event)


class MarkdownView(QTextBrowser):
    def __init__(self, text: str = "", parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self.setObjectName("markdownView")
        self.setFrameShape(QFrame.Shape.NoFrame)
        self.setOpenExternalLinks(True)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.setVerticalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Fixed)
        self.document().documentLayout().documentSizeChanged.connect(
            lambda _size: self._fit())
        # Decoration walks the whole document, so coalesce it while a
        # response is streaming token by token.
        self._decorate_timer = QTimer(self)
        self._decorate_timer.setSingleShot(True)
        self._decorate_timer.setInterval(120)
        self._decorate_timer.timeout.connect(self._decorate)
        self.set_markdown(text)

    def set_markdown(self, text: str) -> None:
        self.setMarkdown(text or "")
        self._decorate_timer.start()
        self._fit()

    def _decorate(self) -> None:
        """Qt's markdown import produces bare code blocks and tables; give
        them the web client's backgrounds, borders, and cell padding."""
        from PySide6.QtGui import (QColor, QTextBlockFormat, QTextCharFormat,
                                   QTextCursor, QTextFormat, QTextTable)
        from .theme import CURRENT
        code_bg = QColor(CURRENT.get("code", "#21211f"))
        border = QColor(CURRENT.get("border", "#43423e"))
        serif = CURRENT.get("serif_font", "Georgia")
        document = self.document()
        frames = [document.rootFrame()]
        while frames:
            frame = frames.pop()
            frames.extend(frame.childFrames())
            if isinstance(frame, QTextTable):
                table_format = frame.format()
                table_format.setCellPadding(5)
                table_format.setCellSpacing(0)
                table_format.setBorder(1.0)
                table_format.setBorderBrush(border)
                table_format.setBorderCollapse(True)
                frame.setFormat(table_format)
        block = document.begin()
        while block.isValid():
            block_format = block.blockFormat()
            fragments = []
            iterator = block.begin()
            while not iterator.atEnd():
                fragments.append(iterator.fragment())
                iterator += 1
            is_code_block = block_format.hasProperty(
                QTextFormat.Property.BlockCodeLanguage) or (
                bool(fragments)
                and all(fragment.charFormat().fontFixedPitch()
                        for fragment in fragments)
                and block.text().strip() != "")
            if block_format.headingLevel() > 0:
                # Headings render in the serif display face, like the
                # Claude client.
                cursor = QTextCursor(block)
                cursor.select(QTextCursor.SelectionType.BlockUnderCursor)
                heading = QTextCharFormat()
                heading.setFontFamilies([serif])
                cursor.mergeCharFormat(heading)
            if is_code_block:
                cursor = QTextCursor(block)
                fill = QTextBlockFormat()
                fill.setBackground(code_bg)
                cursor.mergeBlockFormat(fill)
            else:
                for fragment in fragments:
                    if not fragment.charFormat().fontFixedPitch():
                        continue
                    cursor = QTextCursor(document)
                    cursor.setPosition(fragment.position())
                    cursor.setPosition(fragment.position() + fragment.length(),
                                       QTextCursor.MoveMode.KeepAnchor)
                    inline = QTextCharFormat()
                    inline.setBackground(code_bg)
                    cursor.mergeCharFormat(inline)
            block = block.next()

    def _fit(self) -> None:
        width = max(120, self.viewport().width())
        self.document().setTextWidth(width)
        self.setFixedHeight(max(28, int(self.document().size().height()) + 12))

    def resizeEvent(self, event) -> None:  # noqa: N802
        super().resizeEvent(event)
        QTimer.singleShot(0, self._fit)


class AttachmentTile(QPushButton):
    def __init__(self, metadata: dict, workspace: Path) -> None:
        name = str(metadata.get("name", "attachment"))
        path = Path(name)
        if not path.is_absolute():
            path = workspace / path
        super().__init__(name)
        self.path = path
        self.is_image = metadata.get("kind") == "image"
        self.setObjectName("attachmentTile")
        self.setToolTip(str(path))
        if self.is_image and path.is_file():
            pixmap = QPixmap(str(path))
            if not pixmap.isNull():
                self.setIcon(QIcon(pixmap))
                self.setIconSize(QSize(210, 140))
                self.setMinimumSize(220, 150)
        else:
            self.setIcon(svg_icon("file", 18))
        self.clicked.connect(self.open)

    def open(self) -> None:
        if self.is_image and self.path.is_file():
            ImagePreviewDialog(self.path, self.window()).exec()
        else:
            QDesktopServices.openUrl(QUrl.fromLocalFile(str(self.path)))


class ImagePreviewDialog(QDialog):
    def __init__(self, path: Path, parent: QWidget | None = None) -> None:
        super().__init__(parent)
        self.setWindowTitle(path.name)
        self.resize(1000, 760)
        layout = QVBoxLayout(self)
        bar = QHBoxLayout()
        title = QLabel(path.name)
        title.setObjectName("dialogTitle")
        close = IconButton("close", "Close preview", size=16)
        close.clicked.connect(self.accept)
        bar.addWidget(title, 1)
        bar.addWidget(close)
        layout.addLayout(bar)
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        image = QLabel()
        image.setAlignment(Qt.AlignmentFlag.AlignCenter)
        pixmap = QPixmap(str(path))
        if pixmap.isNull():
            image.setText("This image could not be decoded.")
        else:
            image.setPixmap(pixmap)
            image.setMinimumSize(pixmap.size())
        scroll.setWidget(image)
        layout.addWidget(scroll, 1)


def _diff_html(old: str, new: str, context: int = 2) -> str:
    """GitHub-style colored line diff; overlay colors work on both themes."""
    import difflib
    from html import escape
    added = "rgba(63,185,80,0.20)"
    removed = "rgba(248,81,73,0.22)"
    hunk = "rgba(110,118,129,0.16)"
    rows: list[str] = []
    diff = difflib.unified_diff(
        old.splitlines(), new.splitlines(), lineterm="", n=context)
    for line in diff:
        if line.startswith(("---", "+++")):
            continue
        body = escape(line) or "&nbsp;"
        if line.startswith("@@"):
            rows.append(f'<div style="background:{hunk};opacity:.8">{body}</div>')
        elif line.startswith("+"):
            rows.append(f'<div style="background:{added}">{body}</div>')
        elif line.startswith("-"):
            rows.append(f'<div style="background:{removed}">{body}</div>')
        else:
            rows.append(f"<div>{body}</div>")
    if not rows:
        return ""
    return ('<div style="font-family:\'Cascadia Mono\',Consolas,monospace;'
            'font-size:12px;white-space:pre-wrap">' + "".join(rows) + "</div>")


class _DiffView(QTextBrowser):
    def __init__(self, html: str) -> None:
        super().__init__()
        self.setObjectName("diffView")
        self.setHtml(html)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        # The document reports its pre-layout height here, so size by row
        # count; long diffs scroll vertically inside the card.
        rows = html.count("<div")
        self.setFixedHeight(min(280, max(40, rows * 19 + 18)))


_TOOL_META = {
    "web_search": ("research", "Searched the web", "query"),
    "fetch": ("browser", "Fetched a page", "url"),
    "read_file": ("file", "Read a file", "path"),
    "write_file": ("file", "Wrote a file", "path"),
    "edit_file": ("edit", "Edited a file", "path"),
    "run": ("terminal", "Ran a command", "command"),
    "visual_check": ("image", "Checked the UI visually", "target"),
    "browser": ("browser", "Used the browser", "action"),
    "computer": ("tool", "Controlled the computer", "action"),
    "skill": ("skill", "Loaded a skill", "name"),
    "save_skill": ("skill", "Saved a skill", "name"),
    "remember": ("pin", "Saved a memory", "fact"),
    "history_search": ("chat", "Searched past chats", "query"),
    "mcp": ("tool", "Used an MCP tool", "tool"),
    "subtask": ("code", "Delegated a subtask", "task"),
    "list_dir": ("folder", "Listed a folder", "path"),
    "search": ("research", "Searched the workspace", "text"),
    "todo": ("skill", "Task list", ""),
}


class ToolCard(QFrame):
    def __init__(self, name: str, arguments: str,
                 workspace: Path | None = None) -> None:
        super().__init__()
        self.workspace = workspace or Path.cwd()
        self.tool_name = name
        self.setObjectName("todoCard" if name == "todo" else "toolCard")
        self.outer = QVBoxLayout(self)
        self.outer.setContentsMargins(11, 8, 11, 9)
        self.outer.setSpacing(5)
        icon, title, key_field = _TOOL_META.get(name, ("tool", name, ""))
        try:
            data = json.loads(arguments or "{}")
            data = data if isinstance(data, dict) else {}
        except json.JSONDecodeError:
            data = {}
        header = QHBoxLayout()
        header.setSpacing(8)
        self.toggle = QToolButton()
        self.toggle.setText(title)
        self.toggle.setToolButtonStyle(
            Qt.ToolButtonStyle.ToolButtonTextBesideIcon)
        self.toggle.setCheckable(True)
        self.toggle.setObjectName("toolToggle")
        set_svg_icon(self.toggle, icon, 15, "#d97757")
        header.addWidget(self.toggle)
        detail_text = " ".join(str(data.get(key_field, "")).split())[:80] \
            if key_field else ""
        if not detail_text and name not in {"todo"}:
            detail_text = self._summary(arguments)
        self.detail = QLabel(detail_text)
        self.detail.setObjectName("toolDetail")
        header.addWidget(self.detail, 1)
        self.status = QLabel("Running…")
        self.status.setObjectName("toolStatus")
        header.addWidget(self.status)
        self.outer.addLayout(header)
        if name == "todo":
            self._render_todo(data.get("items"))
        # Edits and new files show their change as a colored diff up front,
        # the way the web client did — not as raw JSON behind a toggle.
        diff_html = self._change_preview(name, arguments)
        if diff_html:
            self.outer.addWidget(_DiffView(diff_html))
        self.result = QPlainTextEdit()
        self.result.setReadOnly(True)
        self.result.setMaximumHeight(220)
        self.result.setVisible(False)
        self.outer.addWidget(self.result)
        if name == "todo":
            self.toggle.setCheckable(False)
        else:
            self.toggle.toggled.connect(self.result.setVisible)

    def _render_todo(self, items) -> None:
        if not isinstance(items, list):
            return
        statuses = {"done": ("todoDone", "check"),
                    "active": ("todoActive", "arrow_right"),
                    "pending": ("todoPending", "circle")}
        for item in items[:12]:
            if not isinstance(item, dict):
                continue
            style, icon = statuses.get(str(item.get("status", "pending")),
                                       statuses["pending"])
            row = QHBoxLayout()
            row.setSpacing(7)
            mark = QLabel()
            mark.setPixmap(svg_pixmap(icon, 12))
            label = QLabel(str(item.get("text", ""))[:160])
            label.setObjectName(style)
            label.setWordWrap(True)
            row.addWidget(mark, 0, Qt.AlignmentFlag.AlignTop)
            row.addWidget(label, 1)
            self.outer.addLayout(row)

    @staticmethod
    def _change_preview(name: str, arguments: str) -> str:
        if name not in {"edit_file", "write_file"}:
            return ""
        try:
            data = json.loads(arguments or "{}")
        except json.JSONDecodeError:
            return ""
        if not isinstance(data, dict):
            return ""
        if name == "edit_file":
            old = str(data.get("old_text", ""))
            new = str(data.get("new_text", ""))
            if not old and not new:
                return ""
            return _diff_html(old, new)
        content = str(data.get("content", ""))
        if not content:
            return ""
        lines = content.splitlines()
        preview = "\n".join(lines[:30])
        if len(lines) > 30:
            preview += f"\n… (+{len(lines) - 30} more lines)"
        return _diff_html("", preview, context=0)

    @staticmethod
    def _summary(arguments: str) -> str:
        try:
            data = json.loads(arguments or "{}")
            if isinstance(data, dict):
                summary = " · ".join(
                    f"{key}: {str(value)[:60]}"
                    for key, value in list(data.items())[:2]
                    if key not in {"content", "old_text", "new_text"})
                if summary:
                    return summary
                return " · ".join(f"{key}: {str(value)[:60]}"
                                  for key, value in list(data.items())[:1])
        except json.JSONDecodeError:
            pass
        return arguments[:100]

    def finish(self, result: str) -> None:
        self.result.setPlainText(result)
        # Size the result box to its content instead of leaving a mostly
        # empty five-line default for one-line results.
        lines = min(11, result.count("\n") + 1)
        metrics = self.result.fontMetrics()
        self.result.setFixedHeight(
            min(220, lines * metrics.lineSpacing() + 30))
        failed = result.lstrip().lower().startswith("error")
        self.status.setText("Failed" if failed else "Done")
        self.status.setProperty("failed", "true" if failed else "false")
        self.status.style().unpolish(self.status)
        self.status.style().polish(self.status)
        if self.tool_name == "todo":
            self.status.setText("")
            return
        self.toggle.setChecked(failed)
        paths = re.findall(
            r"(?:Visual screenshot \[[^\]]+\]|Browser screenshot):\s*(.+?\.(?:png|jpe?g|webp))\s*$",
            result, re.I | re.M)
        if paths:
            gallery = QHBoxLayout()
            for raw in paths[:3]:
                path = Path(raw.strip())
                if not path.is_absolute():
                    path = self.workspace / path
                if path.is_file():
                    gallery.addWidget(AttachmentTile(
                        {"name": str(path), "kind": "image"}, self.workspace))
            gallery.addStretch(1)
            self.outer.addLayout(gallery)


class WelcomePanel(QFrame):
    """Claude-desktop-style empty state: a large serif greeting with the
    composer centered beneath it and a row of prompt pills."""

    prompt_selected = Signal(str)

    _PILLS = {
        "agent": (("Draft a document", "file"), ("Build a spreadsheet", "panels"),
                  ("Make a slide deck", "image"), ("Work on code", "code")),
        "chat": (("Explain a concept", "chat"), ("Improve some writing", "edit"),
                 ("Brainstorm ideas", "skill"), ("Summarize pasted text", "file")),
        "research": (("Compare two products", "research"),
                     ("Market overview", "browser"),
                     ("Evaluate a claim", "research"),
                     ("Technology deep dive", "terminal")),
    }

    def __init__(self) -> None:
        super().__init__()
        self.setObjectName("welcomeHost")
        outer = QVBoxLayout(self)
        outer.setContentsMargins(24, 10, 24, 10)
        outer.setSpacing(0)
        outer.addStretch(3)
        greeting_row = QHBoxLayout()
        greeting_row.addStretch(1)
        self.mark = QLabel()
        self.mark.setPixmap(svg_pixmap("sparkle", 22, "#d97757"))
        self.greeting = QLabel()
        self.greeting.setObjectName("welcomeTitle")
        greeting_row.addWidget(self.mark)
        greeting_row.addSpacing(6)
        greeting_row.addWidget(self.greeting)
        greeting_row.addStretch(1)
        outer.addLayout(greeting_row)
        outer.addSpacing(26)
        self.composer_slot = QVBoxLayout()
        self.composer_slot.setContentsMargins(0, 0, 0, 0)
        composer_row = QHBoxLayout()
        composer_row.addStretch(1)
        composer_row.addLayout(self.composer_slot)
        composer_row.addStretch(1)
        outer.addLayout(composer_row)
        outer.addSpacing(18)
        self.pill_row = QHBoxLayout()
        self.pill_row.setSpacing(8)
        outer.addLayout(self.pill_row)
        outer.addStretch(4)
        self.refresh("agent")

    def refresh(self, mode: str) -> None:
        hour = datetime.now().hour
        self.greeting.setText(
            "Up late?" if hour < 5 else "Good morning" if hour < 12
            else "Good afternoon" if hour < 18 else "Good evening")
        while self.pill_row.count():
            item = self.pill_row.takeAt(0)
            widget = item.widget() if item is not None else None
            if widget is not None:
                widget.deleteLater()
        self.pill_row.addStretch(1)
        for label, icon in self._PILLS.get(mode, self._PILLS["agent"]):
            pill = QPushButton(label)
            pill.setObjectName("promptPill")
            pill.setCursor(Qt.CursorShape.PointingHandCursor)
            set_svg_icon(pill, icon, 13)
            pill.clicked.connect(
                lambda _checked=False, text=label: self.prompt_selected.emit(text))
            self.pill_row.addWidget(pill)
        self.pill_row.addStretch(1)


class TranscriptView(QScrollArea):
    """Conversation timeline with native widgets and incremental updates."""

    revert_requested = Signal(int)
    regenerate_requested = Signal(int)
    edit_requested = Signal(int)

    def __init__(self) -> None:
        super().__init__()
        self.setObjectName("transcript")
        self.setWidgetResizable(True)
        self.setFrameShape(QFrame.Shape.NoFrame)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        # Match the web layout: one centered reading column aligned with the
        # 760px composer (plus this canvas's 24px side margins).
        self.setAlignment(Qt.AlignmentFlag.AlignHCenter | Qt.AlignmentFlag.AlignTop)
        self.viewport().setObjectName("transcriptViewport")
        self.canvas = QWidget()
        self.canvas.setObjectName("transcriptCanvas")
        self.canvas.setMaximumWidth(808)
        self.content_layout = QVBoxLayout(self.canvas)
        self.content_layout.setContentsMargins(24, 30, 24, 44)
        self.content_layout.setSpacing(14)
        self.content_layout.addStretch(1)
        self.setWidget(self.canvas)
        self.workspace = Path.cwd()
        self.mode = "agent"  # set by the window; varies the welcome content
        # Auto-scroll is sticky: it follows generation only while the user
        # is at the bottom. Scrolling up pauses it; returning re-engages.
        self._stick_to_bottom = True
        bar = self.verticalScrollBar()
        bar.sliderMoved.connect(lambda _v: self._update_stick())
        bar.sliderReleased.connect(self._update_stick)
        self._activity_label: QLabel | None = None
        self._activity_base = ""
        self._activity_dots = 0
        self._activity_timer = QTimer(self)
        self._activity_timer.setInterval(450)
        self._activity_timer.timeout.connect(self._tick_activity)
        self._reasoning: MarkdownView | None = None
        self._assistant: MarkdownView | None = None
        self._assistant_text = ""
        self._tool: ToolCard | None = None
        self._activity: QWidget | None = None
        self._scroll_settle = QTimer(self)
        self._scroll_settle.setSingleShot(True)
        self._scroll_settle.setInterval(20)
        self._scroll_settle.timeout.connect(self._scroll_bottom_now)
        self._scroll_anim = QPropertyAnimation(
            self.verticalScrollBar(), b"value", self)
        self._scroll_anim.setDuration(170)
        self._scroll_anim.setEasingCurve(QEasingCurve.Type.OutCubic)
        # Content keeps growing while the animation runs; re-check the real
        # bottom after every animation ends and close any remaining gap.
        self._scroll_anim.finished.connect(self._scroll_settle.start)

    def reset_live(self) -> None:
        self._reasoning = None
        self._assistant = None
        self._assistant_text = ""
        self._tool = None
        self._remove_activity()

    def render_display(self, display: list[dict], workspace: Path) -> None:
        clear_layout(self.content_layout)
        self.workspace = workspace
        last_user = -1
        for display_index, item in enumerate(display):
            if item.get("t") == "user":
                last_user = display_index
            self.add_item(item, display_index=display_index,
                          user_index=last_user)
        self.content_layout.addStretch(1)
        self.reset_live()
        self._stick_to_bottom = True
        self.scroll_bottom()

    def resizeEvent(self, event) -> None:  # noqa: N802
        super().resizeEvent(event)
        # A conversation pinned to its end stays pinned through window
        # resizes; otherwise shrinking the window hides the newest message
        # behind the composer.
        bar = self.verticalScrollBar()
        if bar.maximum() - bar.value() < 400:
            self._scroll_settle.start()

    def _insert(self, widget: QWidget,
                alignment=Qt.AlignmentFlag(0)) -> None:
        # No alignment flag = fill the column width. Passing AlignLeft here
        # (a leftover from the web port) shrank every block to its size hint,
        # which squeezed reports into a narrow strip and pushed user bubbles
        # to the left edge.
        index = max(0, self.content_layout.count() - 1)
        self.content_layout.insertWidget(index, widget, 0, alignment)
        QTimer.singleShot(0, self.scroll_bottom)

    def add_item(self, item: dict, *, display_index: int | None = None,
                 user_index: int | None = None) -> QWidget | None:
        kind = item.get("t")
        text = str(item.get("text", ""))
        if kind == "user":
            outer = QFrame()
            outer.setObjectName("userRow")
            outer_layout = QHBoxLayout(outer)
            outer_layout.setContentsMargins(0, 0, 0, 0)
            outer_layout.setSpacing(6)
            outer_layout.addStretch(1)
            if display_index is not None:
                revert = IconButton("undo", "Revert to this message", size=13)
                revert.setObjectName("messageAction")
                revert.clicked.connect(
                    lambda _checked=False, index=display_index:
                    self.revert_requested.emit(index))
                outer_layout.addWidget(revert)
            box = QFrame()
            box.setObjectName("userBubble")
            layout = QVBoxLayout(box)
            layout.setContentsMargins(14, 9, 14, 10)
            label = QLabel(text)
            label.setWordWrap(True)
            label.setTextInteractionFlags(Qt.TextInteractionFlag.TextSelectableByMouse)
            layout.addWidget(label)
            attachments = item.get("attachments") or []
            if attachments:
                grid = QGridLayout()
                grid.setHorizontalSpacing(8)
                grid.setVerticalSpacing(8)
                for index, metadata in enumerate(attachments[:4]):
                    grid.addWidget(
                        AttachmentTile(metadata, self.workspace),
                        index // 2, index % 2)
                layout.addLayout(grid)
            outer_layout.addWidget(box)
            self._insert(outer)
            return outer
        if kind == "steer":
            # Same row/bubble structure as user messages: a right-aligned
            # label alone gets shrunk to a fragment of its text by the
            # word-wrap size hint.
            outer = QFrame()
            outer_layout = QHBoxLayout(outer)
            outer_layout.setContentsMargins(0, 0, 0, 0)
            outer_layout.addStretch(1)
            box = QFrame()
            box.setObjectName("steerBubble")
            layout = QVBoxLayout(box)
            layout.setContentsMargins(11, 8, 11, 8)
            label = QLabel(f"STEER  {text}")
            label.setWordWrap(True)
            label.setTextInteractionFlags(
                Qt.TextInteractionFlag.TextSelectableByMouse)
            layout.addWidget(label)
            outer_layout.addWidget(box)
            self._insert(outer)
            return label
        if kind == "reasoning":
            details = QFrame()
            details.setObjectName("reasoningCard")
            layout = QVBoxLayout(details)
            toggle = QToolButton()
            toggle.setText("Thought process")
            toggle.setCheckable(True)
            body = MarkdownView(text)
            body.setVisible(False)
            toggle.toggled.connect(body.setVisible)
            layout.addWidget(toggle)
            layout.addWidget(body)
            self._insert(details)
            return body
        if kind == "text":
            outer = QFrame()
            outer.setObjectName("assistantBlock")
            layout = QVBoxLayout(outer)
            layout.setContentsMargins(0, 0, 0, 0)
            layout.setSpacing(2)
            view = MarkdownView(text)
            layout.addWidget(view)
            if user_index is not None and user_index >= 0:
                actions = QHBoxLayout()
                copy_button = IconButton("copy", "Copy response", size=12, text="Copy")
                regenerate = IconButton(
                    "refresh", "Regenerate response", size=12, text="Regenerate")
                edit = IconButton("edit", "Edit prompt", size=12, text="Edit prompt")
                copy_button.setObjectName("messageAction")
                regenerate.setObjectName("messageAction")
                edit.setObjectName("messageAction")
                copy_button.clicked.connect(
                    lambda _checked=False, value=text: self._copy_text(value))
                regenerate.clicked.connect(
                    lambda _checked=False, index=user_index:
                    self.regenerate_requested.emit(index))
                edit.clicked.connect(
                    lambda _checked=False, index=user_index:
                    self.edit_requested.emit(index))
                actions.addWidget(copy_button)
                actions.addWidget(regenerate)
                actions.addWidget(edit)
                actions.addStretch(1)
                layout.addLayout(actions)
            self._insert(outer)
            return view
        if kind == "tool":
            card = ToolCard(str(item.get("name", "tool")),
                            str(item.get("args", "")), self.workspace)
            if item.get("result") is not None:
                card.finish(str(item["result"]))
            self._insert(card)
            return card
        if kind == "skill":
            label = QLabel(f"Skill loaded: {item.get('name', '')}")
            label.setObjectName("notice")
            self._insert(label, Qt.AlignmentFlag.AlignLeft)
            return label
        if kind in {"notice", "error"}:
            label = QLabel(text)
            label.setWordWrap(True)
            label.setObjectName("errorNotice" if kind == "error" else "notice")
            label.setTextInteractionFlags(
                Qt.TextInteractionFlag.TextSelectableByMouse)
            self._insert(label)
            return label
        return None

    @staticmethod
    def _copy_text(text: str) -> None:
        from PySide6.QtWidgets import QApplication
        QApplication.clipboard().setText(text)

    def apply_event(self, event: dict) -> None:
        kind, data = event.get("type"), event.get("data")
        if kind == "user":
            # A new turn begins: drop live references to the previous turn's
            # widgets, or queued follow-ups stream their answer into the old
            # response block ABOVE this user message.
            self.reset_live()
            self._stick_to_bottom = True
            payload = data if isinstance(data, dict) else {"text": str(data)}
            self.add_item({"t": "user", **payload})
        elif kind == "steer":
            self.add_item({"t": "steer", "text": str(data)})
        elif kind == "reasoning_delta":
            self._remove_activity()
            if self._reasoning is None:
                self._reasoning = self.add_item(
                    {"t": "reasoning", "text": ""})  # type: ignore[assignment]
            reasoning = self._reasoning
            if reasoning is not None:
                current = reasoning.toPlainText() + str(data)
                reasoning.set_markdown(current)
        elif kind == "content_delta":
            self._remove_activity()
            if self._assistant is None:
                self._assistant = self.add_item(
                    {"t": "text", "text": ""})  # type: ignore[assignment]
            self._assistant_text += str(data)
            assistant = self._assistant
            if assistant is not None:
                assistant.set_markdown(self._assistant_text)
        elif kind == "tool_call" and isinstance(data, dict):
            self._remove_activity()
            self._tool = ToolCard(str(data.get("name", "tool")),
                                  str(data.get("arguments", "")), self.workspace)
            self._insert(self._tool)
        elif kind == "tool_result" and isinstance(data, dict):
            if self._tool:
                self._tool.finish(str(data.get("result", "")))
        elif kind == "skill_loaded" and isinstance(data, dict):
            self.add_item({"t": "skill", "name": data.get("name", "")})
        elif kind in {"context", "error"}:
            self._remove_activity()
            self.add_item({"t": "error" if kind == "error" else "notice",
                           "text": str(data)})
        elif kind == "activity" and isinstance(data, dict):
            self.show_activity(_phase_label(data))
        elif kind == "heartbeat" and isinstance(data, dict):
            elapsed = int(data.get("elapsed_seconds", 0))
            self.show_activity(f"{_phase_label(data)} · {elapsed}s")
        elif kind == "final":
            self._remove_activity()
            if self._assistant is None and str(data).strip():
                self.add_item({"t": "text", "text": str(data)})
        self.scroll_bottom()

    def show_activity(self, text: str) -> None:
        if self._activity is None:
            frame = QFrame()
            layout = QHBoxLayout(frame)
            layout.setContentsMargins(2, 4, 2, 4)
            layout.setSpacing(8)
            icon = QLabel()
            icon.setPixmap(svg_pixmap("sparkle", 14, "#d97757"))
            effect = QGraphicsOpacityEffect(icon)
            icon.setGraphicsEffect(effect)
            pulse = QPropertyAnimation(effect, b"opacity", icon)
            pulse.setDuration(1100)
            pulse.setStartValue(1.0)
            pulse.setKeyValueAt(0.5, 0.3)
            pulse.setEndValue(1.0)
            pulse.setLoopCount(-1)
            pulse.start(QPropertyAnimation.DeletionPolicy.DeleteWhenStopped)
            label = QLabel()
            label.setObjectName("activity")
            layout.addWidget(icon)
            layout.addWidget(label)
            layout.addStretch(1)
            self._activity = frame
            self._activity_label = label
            self._activity_dots = 0
            self._activity_timer.start()
            self._insert(frame)
        self._activity_base = text.rstrip(".· ")
        self._tick_activity(advance=False)

    def _tick_activity(self, advance: bool = True) -> None:
        if self._activity_label is None:
            return
        if advance:
            self._activity_dots = (self._activity_dots + 1) % 4
        self._activity_label.setText(
            self._activity_base + " " + "·" * self._activity_dots)

    def _remove_activity(self) -> None:
        self._activity_timer.stop()
        self._activity_label = None
        if self._activity is not None:
            self._activity.deleteLater()
            self._activity = None

    def wheelEvent(self, event) -> None:  # noqa: N802
        super().wheelEvent(event)
        self._update_stick()

    def _update_stick(self) -> None:
        bar = self.verticalScrollBar()
        self._stick_to_bottom = bar.value() >= bar.maximum() - 32

    def scroll_bottom(self) -> None:
        self._scroll_bottom_now()
        # Rich-text document height and attachment thumbnails settle on the
        # next layout pass. Coalesce streaming chunks into one settled update
        # so the final lines cannot remain hidden below a stale maximum.
        self._scroll_settle.start()

    def _scroll_bottom_now(self) -> None:
        if not self._stick_to_bottom:
            return
        bar = self.verticalScrollBar()
        delta = bar.maximum() - bar.value()
        if delta <= 0:
            return
        running = (self._scroll_anim.state()
                   == QPropertyAnimation.State.Running)
        if running and self._scroll_anim.endValue() == bar.maximum():
            return
        # Streaming produces many small jumps — snap those; ease big ones.
        # A running animation must be stopped either way or its stale target
        # overwrites the newer position on its next tick.
        self._scroll_anim.stop()
        if delta < 320:
            bar.setValue(bar.maximum())
            return
        self._scroll_anim.setStartValue(bar.value())
        self._scroll_anim.setEndValue(bar.maximum())
        self._scroll_anim.start()


_PHASE_LABELS = {
    "model_wait": "Thinking",
    "continuing": "Continuing response",
    "research_scoping": "Scoping the research",
    "research_searching": "Searching the web",
    "research_reading": "Reading a source",
    "research_reviewing": "Reviewing coverage",
    "research_writing": "Writing the report",
}


def _phase_label(data: dict) -> str:
    raw = str(data.get("phase", "working"))
    label = _PHASE_LABELS.get(raw, raw.replace("_", " ").title())
    if raw == "research_searching" and data.get("round"):
        label += f" (round {data.get('round')}/{data.get('rounds', '?')})"
    return label


class JobWatcher(QThread):
    event_received = Signal(dict)
    stream_finished = Signal(str)

    def __init__(self, job) -> None:
        super().__init__()
        self.job = job

    def run(self) -> None:
        state = "done"
        while not self.isInterruptionRequested():
            try:
                item = self.job.events.get(timeout=0.25)
            except queue.Empty:
                self.event_received.emit({
                    "type": "heartbeat", "data": self.job.activity_snapshot()})
                continue
            if item is None:
                state = self.job.state
                break
            self.event_received.emit(item)
        self.stream_finished.emit(state)


class SessionList(QListWidget):
    delete_requested = Signal(str)
    rename_requested = Signal(str)

    def __init__(self) -> None:
        super().__init__()
        self.setObjectName("sessionList")
        self.setSelectionMode(QAbstractItemView.SelectionMode.SingleSelection)
        self.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        # Long titles must elide, not grow a horizontal scrollbar under the
        # sidebar (a leftover from the CSS text-overflow behavior).
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.setTextElideMode(Qt.TextElideMode.ElideRight)
        self.setWordWrap(False)

    @staticmethod
    def _relative_time(timestamp: float) -> str:
        delta = max(0, int(time.time() - timestamp))
        if delta < 60:
            return "just now"
        if delta < 3600:
            return f"{delta // 60}m ago"
        if delta < 86400:
            return f"{delta // 3600}h ago"
        if delta < 7 * 86400:
            return f"{delta // 86400}d ago"
        return datetime.fromtimestamp(timestamp).strftime("%b %d")

    def set_sessions(self, sessions: list[dict], current_id: str | None) -> None:
        self.clear()
        for session in sessions:
            item = QListWidgetItem()
            item.setData(Qt.ItemDataRole.UserRole, session["id"])
            row = QWidget()
            layout = QVBoxLayout(row)
            layout.setContentsMargins(10, 6, 8, 7)
            layout.setSpacing(1)
            title_row = QHBoxLayout()
            title_row.setSpacing(5)
            title = QLabel(session.get("title") or "New chat")
            title.setObjectName("sessionTitle")
            metrics = title.fontMetrics()
            title.setText(metrics.elidedText(
                title.text(), Qt.TextElideMode.ElideRight, 180))
            if session.get("pinned"):
                pin = QLabel()
                pin.setPixmap(svg_pixmap("pin", 11))
                title_row.addWidget(pin)
            title_row.addWidget(title, 1)
            layout.addLayout(title_row)
            if session.get("running"):
                state = "working…"
            elif session.get("queued"):
                state = "queued"
            elif session.get("pending_count"):
                state = f"{session['pending_count']} queued"
            else:
                state = ""
            when = self._relative_time(float(session.get("updated") or 0))
            subtitle = QLabel(f"{when} · {state}" if state else when)
            subtitle.setObjectName(
                "sessionStateBusy" if state else "sessionSubtitle")
            layout.addWidget(subtitle)
            item.setSizeHint(row.sizeHint())
            self.addItem(item)
            self.setItemWidget(item, row)
            if session["id"] == current_id:
                self.setCurrentItem(item)


class WheelGuard(QObject):
    """Prevent accidental value changes when scrolling a settings dialog."""

    def eventFilter(self, watched, event) -> bool:  # noqa: N802
        if event.type() == QEvent.Type.Wheel and not watched.hasFocus():
            event.ignore()
            return True
        return super().eventFilter(watched, event)
