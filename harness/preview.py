"""Server-side artifact previews: render workspace files to clean HTML.

This is what lets the web UI show Word docs, spreadsheets, decks, PDFs,
markdown and code as live artifacts — not just raw HTML snippets.
"""
from __future__ import annotations

import html
import json
import re
from pathlib import Path

BASE_CSS = """
/* previews render like paper: always light, independent of app theme */
:root { color-scheme: light; }
* { box-sizing: border-box; }
body { font: 15px/1.65 -apple-system, "Segoe UI", system-ui, sans-serif;
  color: #1f1e1b; background: #fffdf9; margin: 0; padding: 32px 40px;
  max-width: 860px; margin-inline: auto; }
h1, h2, h3, h4 { font-family: Georgia, serif; line-height: 1.3; margin: 1.2em 0 .4em; }
h1 { font-size: 26px; } h2 { font-size: 21px; } h3 { font-size: 17px; }
table { border-collapse: collapse; margin: 14px 0; font-size: 13.5px; }
th, td { border: 1px solid #d8d5cc; padding: 6px 12px; text-align: left; }
th { background: #f0eee7; font-weight: 600; }
pre { background: #f4f2ec; padding: 14px 16px; border-radius: 10px;
  overflow-x: auto; font: 13px/1.5 Consolas, monospace; white-space: pre-wrap; }
code { background: #f4f2ec; padding: 1px 5px; border-radius: 4px;
  font-family: Consolas, monospace; font-size: 13px; }
blockquote { border-left: 3px solid #c96442; margin: 12px 0; padding: 2px 16px; color: #75726a; }
hr { border: none; border-top: 1px solid #d8d5cc; margin: 18px 0; }
.slide { border: 1px solid #d8d5cc; border-radius: 14px; padding: 22px 28px;
  margin: 18px 0; box-shadow: 0 1px 4px rgba(0,0,0,.06); }
.slide h2 { margin-top: 0; color: #c96442; }
.slide .notes { font-size: 12.5px; color: #75726a; border-top: 1px dashed #d8d5cc;
  margin-top: 12px; padding-top: 8px; font-style: italic; }
.sheetname { color: #c96442; }
.page-label { font-size: 11px; color: #a5a297; text-transform: uppercase;
  letter-spacing: .8px; margin: 22px 0 6px; }
.meta { color:#a5a297; font-size: 12px; margin-bottom: 18px; }
"""


def _page(title: str, body: str) -> str:
    return (f"<!doctype html><html><head><meta charset='utf-8'>"
            f"<title>{html.escape(title)}</title><style>{BASE_CSS}</style>"
            f"</head><body>{body}</body></html>")


def _esc(s) -> str:
    return html.escape(str(s if s is not None else ""))


# ---------- markdown (small, matches what the skills emit) ----------
def _md_inline(s: str) -> str:
    s = _esc(s)
    s = re.sub(r"`([^`]+)`", r"<code>\1</code>", s)
    s = re.sub(r"\*\*([^*]+)\*\*", r"<b>\1</b>", s)
    s = re.sub(r"(?<![\w*])\*([^*\n]+)\*(?![\w*])", r"<i>\1</i>", s)
    s = re.sub(r"\[([^\]]+)\]\((https?:[^)\s]+)\)",
               r'<a href="\2" target="_blank">\1</a>', s)
    return s


def md_to_html(src: str) -> str:
    out, lines, i = [], src.split("\n"), 0
    while i < len(lines):
        line = lines[i]
        if line.strip().startswith("```"):
            code = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            out.append(f"<pre>{_esc(chr(10).join(code))}</pre>")
            continue
        if not line.strip():
            i += 1
            continue
        m = re.match(r"^(#{1,4})\s+(.*)", line)
        if m:
            n = len(m.group(1))
            out.append(f"<h{n}>{_md_inline(m.group(2))}</h{n}>")
            i += 1
            continue
        if re.match(r"^(-{3,}|\*{3,})\s*$", line.strip()):
            out.append("<hr>")
            i += 1
            continue
        if line.strip().startswith(">"):
            q = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                q.append(lines[i].strip().lstrip("> "))
                i += 1
            out.append(f"<blockquote>{_md_inline(' '.join(q))}</blockquote>")
            continue
        if re.match(r"^\|.*\|", line) and i + 1 < len(lines) \
                and re.match(r"^\|[\s:|-]+\|", lines[i + 1]):
            cells = lambda r: [c.strip() for c in r.strip().strip("|").split("|")]
            rows_html = "<tr>" + "".join(
                f"<th>{_md_inline(c)}</th>" for c in cells(line)) + "</tr>"
            i += 2
            while i < len(lines) and re.match(r"^\|.*\|", lines[i]):
                rows_html += "<tr>" + "".join(
                    f"<td>{_md_inline(c)}</td>" for c in cells(lines[i])) + "</tr>"
                i += 1
            out.append(f"<table>{rows_html}</table>")
            continue
        m = re.match(r"^(\s*)([-*]|\d+[.)])\s+(.*)", line)
        if m:
            tag = "ol" if m.group(2)[0].isdigit() else "ul"
            items = []
            while i < len(lines):
                m2 = re.match(r"^(\s*)([-*]|\d+[.)])\s+(.*)", lines[i])
                if not m2:
                    break
                items.append(f"<li>{_md_inline(m2.group(3))}</li>")
                i += 1
            out.append(f"<{tag}>{''.join(items)}</{tag}>")
            continue
        para = [line]
        i += 1
        while i < len(lines) and lines[i].strip() and not re.match(
                r"^(#{1,4}\s|```|>|\||(\s*([-*]|\d+[.)])\s))", lines[i]):
            para.append(lines[i])
            i += 1
        out.append(f"<p>{_md_inline(' '.join(p.strip() for p in para))}</p>")
    return "".join(out)


# ---------- office renderers ----------
def _docx_html(path: Path) -> str:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    doc = Document(str(path))
    parts = []
    for child in doc.element.body.iterchildren():
        if child.tag.endswith("}p"):
            p = Paragraph(child, doc)
            text = _esc(p.text.strip())
            if not text:
                continue
            style = (p.style.name or "").lower()
            if style.startswith("heading"):
                n = min(int("".join(c for c in style if c.isdigit()) or 1), 4)
                parts.append(f"<h{n}>{text}</h{n}>")
            elif "quote" in style:
                parts.append(f"<blockquote>{text}</blockquote>")
            elif "list" in style:
                parts.append(f"<ul><li>{text}</li></ul>")
            else:
                parts.append(f"<p>{text}</p>")
        elif child.tag.endswith("}tbl"):
            t = Table(child, doc)
            rows = []
            for ri, row in enumerate(t.rows):
                tag = "th" if ri == 0 else "td"
                rows.append("<tr>" + "".join(
                    f"<{tag}>{_esc(c.text.strip())}</{tag}>"
                    for c in row.cells) + "</tr>")
            parts.append(f"<table>{''.join(rows)}</table>")
    return "".join(parts) or "<p class='meta'>(empty document)</p>"


def _xlsx_html(path: Path) -> str:
    from openpyxl import load_workbook
    parts = []
    # values pass (computed where cached) + formulas pass for tooltips
    wb_val = load_workbook(str(path), data_only=True)
    wb_f = load_workbook(str(path), data_only=False)
    for name in wb_val.sheetnames:
        wsv, wsf = wb_val[name], wb_f[name]
        parts.append(f"<h2 class='sheetname'>{_esc(name)}</h2>")
        rows = []
        for ri, (rv, rf) in enumerate(zip(
                wsv.iter_rows(values_only=True),
                wsf.iter_rows(values_only=True)), 1):
            if ri > 200:
                parts.append("<p class='meta'>…more rows not shown</p>")
                break
            tag = "th" if ri == 1 else "td"
            cells = []
            for v, f in zip(rv, rf):
                shown = v if v is not None else (f if f is not None else "")
                tip = (f" title='{_esc(f)}'"
                       if isinstance(f, str) and f.startswith("=") else "")
                cells.append(f"<{tag}{tip}>{_esc(shown)}</{tag}>")
            rows.append("<tr>" + "".join(cells) + "</tr>")
        parts.append(f"<table>{''.join(rows)}</table>")
        parts.append("<p class='meta'>hover a cell to see its formula</p>")
    return "".join(parts)


def _pptx_html(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        body = [f"<div class='page-label'>slide {i}</div>", "<div class='slide'>"]
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if t:
                    texts.append((para.level, t))
        if texts:
            first = texts[0][1]
            body.append(f"<h2>{_esc(first)}</h2>")
            rest = texts[1:]
            if rest:
                body.append("<ul>" + "".join(
                    f"<li style='margin-left:{lvl * 18}px'>{_esc(t)}</li>"
                    for lvl, t in rest) + "</ul>")
        if slide.has_notes_slide:
            n = slide.notes_slide.notes_text_frame.text.strip()
            if n:
                body.append(f"<div class='notes'>🗒 {_esc(n)}</div>")
        body.append("</div>")
        parts.append("".join(body))
    return "".join(parts) or "<p class='meta'>(empty deck)</p>"


def _pdf_html(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages[:20], 1):
        parts.append(f"<div class='page-label'>page {i} of {len(reader.pages)}</div>")
        parts.append(f"<pre>{_esc(page.extract_text() or '(no extractable text)')}</pre>")
    if len(reader.pages) > 20:
        parts.append("<p class='meta'>…more pages not shown</p>")
    return "".join(parts)


def _csv_html(path: Path) -> str:
    import csv
    rows = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for ri, row in enumerate(csv.reader(f), 1):
            if ri > 300:
                rows.append("<tr><td colspan='99' class='meta'>…more rows</td></tr>")
                break
            tag = "th" if ri == 1 else "td"
            rows.append("<tr>" + "".join(
                f"<{tag}>{_esc(c)}</{tag}>" for c in row) + "</tr>")
    return f"<table>{''.join(rows)}</table>"


CODE_EXTS = {".py", ".js", ".ts", ".css", ".ps1", ".sh", ".bat", ".c", ".cpp",
             ".java", ".rs", ".go", ".toml", ".yaml", ".yml", ".xml", ".sql"}


def build_preview(path: Path) -> str:
    """Return a complete self-contained HTML document previewing the file."""
    ext = path.suffix.lower()
    title = path.name
    try:
        if ext == ".docx":
            return _page(title, _docx_html(path))
        if ext == ".xlsx":
            return _page(title, _xlsx_html(path))
        if ext == ".pptx":
            return _page(title, _pptx_html(path))
        if ext == ".pdf":
            return _page(title, _pdf_html(path))
        if ext == ".csv":
            return _page(title, _csv_html(path))
        if ext in (".md", ".markdown"):
            return _page(title, md_to_html(
                path.read_text(encoding="utf-8", errors="replace")))
        if ext == ".json":
            data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
            return _page(title, f"<pre>{_esc(json.dumps(data, indent=2, ensure_ascii=False))}</pre>")
        if ext in CODE_EXTS or ext in (".txt", ".log"):
            return _page(title, f"<pre>{_esc(path.read_text(encoding='utf-8', errors='replace')[:60000])}</pre>")
    except Exception as e:
        return _page(title, f"<p>Could not render preview: {_esc(e)}</p>")
    return _page(title, "<p class='meta'>No preview available for this file type — use download.</p>")


MEDIA_EXTS = {".mp4", ".webm", ".mov", ".mp3", ".wav", ".ogg"}

PREVIEWABLE = {".docx", ".xlsx", ".pptx", ".pdf", ".csv", ".md", ".markdown",
               ".json", ".txt", ".log", ".html", ".htm", ".svg", ".png",
               ".jpg", ".jpeg", ".gif", ".webp"} | CODE_EXTS | MEDIA_EXTS
