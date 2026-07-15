"""Dump .pptx contents as a Markdown-ish outline.

Usage: python read_pptx.py file.pptx
"""
import sys

# Windows pipes default to a legacy code page; documents contain
# arbitrary Unicode, so never let printing crash the read.
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")


from pptx import Presentation

MAX_CHARS = 8000


def main(path):
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if text:
                    out.append(("  " * para.level) + "- " + text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"Notes: {notes}")
        out.append("")
    text = "\n".join(out)
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + "\n...[truncated]"
    print(text or "(empty presentation)")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit("Usage: python read_pptx.py file.pptx")
    main(sys.argv[1])
