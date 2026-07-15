"""Convert a Markdown outline into a designed .pptx deck.

Usage: python md2pptx.py outline.md output.pptx

The outline is content-only; this script owns the design system: themed
palettes, a title/section/content/stat/quote/table/columns/closing slide
family, consistent typography, accent geometry, and footers. See the
presentations skill for the outline syntax.
"""
import datetime
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

XML_INVALID_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")
MAX_SLIDES = 60

# Each theme: page background, ink (body text), muted, accent, accent_dark
# (title/section backgrounds), card (stat/table band fill), on_dark text.
THEMES = {
    "ink": dict(page="F7F5F1", ink="1F1E1B", muted="6E6A61", accent="C96442",
                dark="26241F", card="EDE9E1", on_dark="F5F2EC"),
    "ocean": dict(page="F4F7FA", ink="16232E", muted="5C6B7A", accent="1273B5",
                  dark="0E2438", card="E3ECF4", on_dark="EFF5FA"),
    "forest": dict(page="F5F7F3", ink="1D261C", muted="60705E", accent="3E7C4F",
                   dark="15281B", card="E4EBE0", on_dark="F0F5EE"),
    "slate": dict(page="F6F6F8", ink="1E2024", muted="63676E", accent="5B5FC7",
                  dark="1B1D2A", card="E7E7EE", on_dark="F2F2F7"),
    "sunset": dict(page="FBF6F1", ink="2A211C", muted="77685C", accent="D4622A",
                   dark="332014", card="F3E8DE", on_dark="FAF1EA"),
}
TITLE_FONT = "Georgia"
BODY_FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)


def clean(text):
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text).strip()
    return XML_INVALID_RE.sub("", text)


def rgb(hex_str):
    return RGBColor.from_string(hex_str)


# ---------------- outline parsing ----------------

KIND_RE = re.compile(r"^\[(section|stats|quote|table|cols|closing)\]\s*(.*)$",
                     re.IGNORECASE)


def parse(md_path):
    lines = Path(md_path).read_text(encoding="utf-8").splitlines()
    deck = {"title": "", "subtitle": "", "theme": "ink", "footer": "",
            "date": "", "slides": []}
    slide = None
    for line in lines:
        stripped = line.strip()
        if not stripped:
            if slide is not None:
                slide["blocks"].append(("blank", ""))
            continue
        lower = stripped.lower()
        if slide is None and lower.startswith("theme:"):
            requested = stripped.split(":", 1)[1].strip().lower()
            if requested in THEMES:
                deck["theme"] = requested
            continue
        if slide is None and lower.startswith("footer:"):
            deck["footer"] = clean(stripped.split(":", 1)[1])
            continue
        if slide is None and lower.startswith("date:"):
            deck["date"] = clean(stripped.split(":", 1)[1])
            continue
        if stripped.startswith("# ") and not deck["title"]:
            deck["title"] = clean(stripped[2:])
            continue
        if stripped.startswith("## "):
            body = stripped[3:].strip()
            kind, title = "bullets", body
            m = KIND_RE.match(body)
            if m:
                kind, title = m.group(1).lower(), clean(m.group(2))
            slide = {"kind": kind, "title": clean(title), "kicker": "",
                     "blocks": [], "notes": []}
            deck["slides"].append(slide)
            continue
        if slide is None:
            if deck["title"] and not deck["subtitle"]:
                deck["subtitle"] = clean(stripped)
            continue
        if lower.startswith("notes:"):
            slide["notes"].append(clean(stripped[6:]))
            continue
        if lower.startswith("kicker:"):
            slide["kicker"] = clean(stripped.split(":", 1)[1]).upper()
            continue
        if stripped == "---":
            slide["blocks"].append(("split", ""))
            continue
        if stripped.startswith("|"):
            slide["blocks"].append(("row", stripped))
            continue
        if stripped.startswith(">"):
            slide["blocks"].append(("quote", clean(stripped.lstrip("> "))))
            continue
        if stripped.startswith("—") or stripped.startswith("--"):
            slide["blocks"].append(("attrib", clean(stripped.lstrip("—- "))))
            continue
        m = re.match(r"^(\s*)[-*]\s+(.*)$", line)
        if m:
            slide["blocks"].append(
                ("bullet", (len(m.group(1)) // 2, clean(m.group(2)))))
            continue
        slide["blocks"].append(("para", clean(stripped)))
    return deck


# ---------------- drawing primitives ----------------

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, color, rounded=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    if rounded:
        try:
            shape.adjustments[0] = 0.06
        except (IndexError, ValueError):
            pass
    solid(shape, color)
    return shape


def circle(slide, x, y, d, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    solid(shape, color)
    return shape


def text(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def para(frame, content, *, size, color, first=False, bold=False,
         italic=False, font=BODY_FONT, align=PP_ALIGN.LEFT, before=0,
         after=0, caps=False):
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.text = content.upper() if caps else content
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    return p


def set_background(slide, color):
    rect(slide, 0, 0, W, H, color)


def decor_circles(slide, theme, cx, cy):
    """Overlapping translucent-feel circles (solid tints on dark)."""
    accent = rgb(theme["accent"])
    dark_mix = RGBColor((accent[0] + int(theme["dark"][0:2], 16)) // 2,
                        (accent[1] + int(theme["dark"][2:4], 16)) // 2,
                        (accent[2] + int(theme["dark"][4:6], 16)) // 2)
    circle(slide, cx, cy, Inches(3.2), dark_mix)
    circle(slide, cx + Inches(1.9), cy + Inches(1.1), Inches(1.7), accent)


def footer_row(slide, theme, deck, number):
    if deck["footer"]:
        frame = text(slide, MARGIN, H - Inches(0.5), Inches(6), Inches(0.35))
        para(frame, deck["footer"], size=10.5, color=rgb(theme["muted"]),
             first=True)
    frame = text(slide, W - Inches(1.6), H - Inches(0.5), Inches(0.85),
                 Inches(0.35))
    para(frame, f"{number:02d}", size=10.5, color=rgb(theme["muted"]),
         first=True, align=PP_ALIGN.RIGHT)


def content_header(slide, theme, sl):
    y = Inches(0.55)
    if sl["kicker"]:
        frame = text(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.35))
        para(frame, sl["kicker"], size=12, color=rgb(theme["accent"]),
             first=True, bold=True, caps=True)
        y += Inches(0.38)
    frame = text(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.85))
    para(frame, sl["title"], size=29, color=rgb(theme["ink"]), first=True,
         bold=True, font=TITLE_FONT)
    rect(slide, MARGIN, y + Inches(0.78), Inches(0.85), Pt(3.2),
         rgb(theme["accent"]))
    return y + Inches(1.1)


# ---------------- slide builders ----------------

def build_title(prs, deck, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["dark"]))
    decor_circles(slide, theme, W - Inches(3.4), H - Inches(2.6))
    rect(slide, MARGIN, Inches(2.18), Inches(0.9), Pt(4), rgb(theme["accent"]))
    frame = text(slide, MARGIN, Inches(2.45), W - Inches(4.2), Inches(2.6))
    para(frame, deck["title"], size=46, color=rgb(theme["on_dark"]),
         first=True, bold=True, font=TITLE_FONT)
    if deck["subtitle"]:
        para(frame, deck["subtitle"], size=19, color=rgb(theme["muted"]),
             before=14)
    stamp = deck["date"] or datetime.date.today().strftime("%B %Y")
    frame = text(slide, MARGIN, H - Inches(0.85), Inches(7), Inches(0.4))
    line = " · ".join(part for part in (deck["footer"], stamp) if part)
    para(frame, line, size=12, color=rgb(theme["muted"]), first=True)


def build_section(prs, deck, theme, sl, number, section_index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["dark"]))
    rect(slide, 0, 0, Inches(0.22), H, rgb(theme["accent"]))
    frame = text(slide, MARGIN, Inches(2.15), Inches(2.5), Inches(1.2))
    para(frame, f"{section_index:02d}", size=54, color=rgb(theme["accent"]),
         first=True, bold=True, font=TITLE_FONT)
    frame = text(slide, MARGIN, Inches(3.35), W - 2 * MARGIN, Inches(1.8))
    para(frame, sl["title"], size=38, color=rgb(theme["on_dark"]),
         first=True, bold=True, font=TITLE_FONT)
    paras = [b for kind, b in sl["blocks"] if kind == "para"]
    if paras:
        para(frame, paras[0], size=16, color=rgb(theme["muted"]), before=12)
    footer_row(slide, theme, deck, number)


def build_bullets(prs, deck, theme, sl, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["page"]))
    top = content_header(slide, theme, sl)
    bullets = [b for kind, b in sl["blocks"] if kind == "bullet"]
    paras = [b for kind, b in sl["blocks"] if kind == "para"]
    if bullets:
        count = len(bullets)
        size = 20 if count <= 5 else 17 if count <= 7 else 15
        gap = 14 if count <= 5 else 9
        frame = text(slide, MARGIN + Inches(0.05), top + Inches(0.25),
                     W - 2 * MARGIN - Inches(0.1), H - top - Inches(1.1))
        first = True
        for level, content in bullets[:10]:
            p = para(frame, "", size=size if level == 0 else size - 2.5,
                     color=rgb(theme["ink"] if level == 0 else theme["muted"]),
                     first=first, after=gap)
            first = False
            p.level = min(level, 3)
            marker = p.runs[0] if p.runs else p.add_run()
            marker.text = ("▪  " if level == 0 else "–  ")
            marker.font.color.rgb = rgb(theme["accent"])
            marker.font.bold = True
            body = p.add_run()
            body.text = content
            body.font.color.rgb = rgb(
                theme["ink"] if level == 0 else theme["muted"])
    elif paras:
        frame = text(slide, Inches(2.1), Inches(2.6), W - Inches(4.2),
                     Inches(2.6), anchor=MSO_ANCHOR.MIDDLE)
        first = True
        for content in paras[:3]:
            para(frame, content, size=27, color=rgb(theme["ink"]),
                 first=first, align=PP_ALIGN.CENTER, font=TITLE_FONT,
                 after=10)
            first = False
    footer_row(slide, theme, deck, number)


def build_stats(prs, deck, theme, sl, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["page"]))
    top = content_header(slide, theme, sl)
    entries = []
    for kind, payload in sl["blocks"]:
        if kind != "bullet":
            continue
        _level, content = payload
        value, _, label = content.partition("|")
        entries.append((value.strip(), label.strip()))
    entries = entries[:4] or [("—", "add '- value | label' lines")]
    gap = Inches(0.3)
    width = Emu(int((W - 2 * MARGIN - gap * (len(entries) - 1))
                    / len(entries)))
    card_top = top + Inches(0.55)
    card_height = Inches(2.9)
    for index, (value, label) in enumerate(entries):
        x = MARGIN + index * (width + gap)
        rect(slide, x, card_top, width, card_height, rgb(theme["card"]),
             rounded=True)
        rect(slide, x, card_top, width, Pt(4), rgb(theme["accent"]))
        frame = text(slide, x + Inches(0.25), card_top + Inches(0.55),
                     width - Inches(0.5), card_height - Inches(0.9))
        para(frame, value, size=40, color=rgb(theme["accent"]), first=True,
             bold=True, font=TITLE_FONT)
        para(frame, label, size=13.5, color=rgb(theme["ink"]), before=10)
    footer_row(slide, theme, deck, number)


def build_quote(prs, deck, theme, sl, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["page"]))
    quotes = [b for kind, b in sl["blocks"] if kind in {"quote", "para"}]
    attribution = next(
        (b for kind, b in sl["blocks"] if kind == "attrib"), "")
    frame = text(slide, Inches(1.1), Inches(0.9), Inches(2), Inches(2))
    para(frame, "“", size=130, color=rgb(theme["accent"]), first=True,
         bold=True, font=TITLE_FONT)
    frame = text(slide, Inches(2.0), Inches(2.35), W - Inches(4.0),
                 Inches(3.0), anchor=MSO_ANCHOR.MIDDLE)
    first = True
    for content in quotes[:2]:
        para(frame, content, size=26, color=rgb(theme["ink"]), first=first,
             italic=True, font=TITLE_FONT, after=8)
        first = False
    if attribution:
        para(frame, "— " + attribution, size=15, color=rgb(theme["muted"]),
             before=16)
    footer_row(slide, theme, deck, number)


def build_table(prs, deck, theme, sl, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["page"]))
    top = content_header(slide, theme, sl)
    raw = [payload for kind, payload in sl["blocks"] if kind == "row"]
    rows = []
    for line in raw:
        if re.match(r"^\s*\|?[\s:|-]+\|[\s:|-]*$", line) and "-" in line:
            continue
        rows.append([clean(cell) for cell in
                     line.strip().strip("|").split("|")])
    if not rows:
        rows = [["Add a markdown table under this heading"]]
    cols = max(len(row) for row in rows)
    rows = [row + [""] * (cols - len(row)) for row in rows[:12]]
    height = min(Inches(0.52) * len(rows), H - top - Inches(1.0))
    shape = slide.shapes.add_table(
        len(rows), cols, MARGIN, top + Inches(0.35), W - 2 * MARGIN, height)
    table = shape.table
    # replace the default banded style with the theme's colors
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # "no style"
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.margin_left = cell.margin_right = Inches(0.12)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if row_index == 0:
                cell.fill.fore_color.rgb = rgb(theme["dark"])
            elif row_index % 2 == 1:
                cell.fill.fore_color.rgb = rgb(theme["page"])
            else:
                cell.fill.fore_color.rgb = rgb(theme["card"])
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = value
            run.font.size = Pt(13)
            run.font.name = BODY_FONT
            run.font.bold = row_index == 0
            run.font.color.rgb = rgb(
                theme["on_dark"] if row_index == 0 else theme["ink"])
    footer_row(slide, theme, deck, number)


def build_cols(prs, deck, theme, sl, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["page"]))
    top = content_header(slide, theme, sl)
    columns = [[]]
    for kind, payload in sl["blocks"]:
        if kind == "split":
            columns.append([])
        elif kind in {"bullet", "para"}:
            columns[-1].append((kind, payload))
    columns = [col for col in columns if col][:2] or [[]]
    width = Emu(int((W - 2 * MARGIN - Inches(0.5)) / max(len(columns), 1)))
    for index, column in enumerate(columns):
        x = MARGIN + index * (width + Inches(0.5))
        rect(slide, x, top + Inches(0.3), width, H - top - Inches(1.25),
             rgb(theme["card"]), rounded=True)
        frame = text(slide, x + Inches(0.3), top + Inches(0.6),
                     width - Inches(0.6), H - top - Inches(1.9))
        first = True
        for kind, payload in column[:8]:
            if kind == "para":
                para(frame, payload, size=15.5, color=rgb(theme["accent"]),
                     first=first, bold=True, after=8)
            else:
                _level, content = payload
                p = para(frame, "", size=14.5, color=rgb(theme["ink"]),
                         first=first, after=7)
                marker = p.runs[0] if p.runs else p.add_run()
                marker.text = "▪  "
                marker.font.color.rgb = rgb(theme["accent"])
                body = p.add_run()
                body.text = content
                body.font.color.rgb = rgb(theme["ink"])
            first = False
    footer_row(slide, theme, deck, number)


def build_closing(prs, deck, theme, sl, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, rgb(theme["dark"]))
    decor_circles(slide, theme, Inches(-1.4), H - Inches(2.2))
    frame = text(slide, Inches(2.0), Inches(2.7), W - Inches(4.0),
                 Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
    para(frame, sl["title"] or "Thank you", size=44,
         color=rgb(theme["on_dark"]), first=True, bold=True, font=TITLE_FONT,
         align=PP_ALIGN.CENTER)
    lines = [b for kind, b in sl["blocks"] if kind in {"para", "quote"}]
    for content in lines[:3]:
        para(frame, content, size=16, color=rgb(theme["muted"]), before=10,
             align=PP_ALIGN.CENTER)


BUILDERS = {"stats": build_stats, "quote": build_quote, "table": build_table,
            "cols": build_cols}


def build(md_path, out_path):
    deck = parse(md_path)
    if len(deck["slides"]) > MAX_SLIDES:
        sys.exit(f"Error: at most {MAX_SLIDES} slides are supported")
    theme = THEMES[deck["theme"]]
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    number = 0
    if deck["title"]:
        build_title(prs, deck, theme)
        number += 1
    section_index = 0
    for sl in deck["slides"]:
        number += 1
        if sl["kind"] == "section":
            section_index += 1
            build_section(prs, deck, theme, sl, number, section_index)
        elif sl["kind"] == "closing":
            build_closing(prs, deck, theme, sl, number)
        else:
            BUILDERS.get(sl["kind"], build_bullets)(
                prs, deck, theme, sl, number)
        if sl["notes"]:
            prs.slides[-1].notes_slide.notes_text_frame.text = \
                " ".join(sl["notes"])
    output = Path(out_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Saved {out_path} ({len(prs.slides)} slides, "
          f"theme '{deck['theme']}')")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        sys.exit("Usage: python md2pptx.py outline.md output.pptx")
    build(sys.argv[1], sys.argv[2])
